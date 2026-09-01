import tempfile
import threading
import types
import unittest
import zipfile
import xml.etree.ElementTree as ET
from io import BytesIO
from pathlib import Path
from unittest.mock import Mock, patch

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from actions import presentation_maker as maker
from actions.presentations import assets as presentation_assets
from actions.presentations.models import PresentationRequest, SourceBundle
from actions.presentations.models3d import (
    Native3DAsset,
    inspect_native_3d,
)
from actions.presentations.quality import verify_pptx
from actions.presentations.sources import ingest_sources
from agent import planner
from agent.task_queue import TaskQueue
from core.qa_audit import declared_tools


def _sample_plan():
    slide_types = ["cover", "statement", "two_column", "process", "comparison", "closing"]
    return {
        "title": "The Future of Clean Energy",
        "subtitle": "A practical transition plan",
        "slides": [
            {
                "type": kind,
                "kicker": "ENERGY SHIFT",
                "title": f"Claim {index} moves the story forward",
                "subtitle": "Concise support for the claim.",
                "bullets": ["First supporting idea", "Second supporting idea"],
                "visual_label": "HOW IT WORKS",
                "visual_items": ["Frame", "Build", "Scale"],
                "left_label": "TODAY",
                "right_label": "NEXT",
                "left_items": ["Fragmented", "Costly"],
                "right_items": ["Connected", "Efficient"],
                "source_note": "",
            }
            for index, kind in enumerate(slide_types, start=1)
        ],
    }


_MODEL_FRAME_XML = b'''<p:graphicFrame
 xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
 xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
 xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
 xmlns:am3d="http://schemas.microsoft.com/office/drawing/2017/model3d">
 <p:nvGraphicFramePr><p:cNvPr id="2" name="3D Model 1" descr="Test Cell"/><p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>
 <p:xfrm><a:off x="0" y="0"/><a:ext cx="1000000" cy="1000000"/></p:xfrm>
 <a:graphic><a:graphicData uri="http://schemas.microsoft.com/office/drawing/2017/model3d">
  <am3d:model3d r:embed="rId2">
   <am3d:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="1000000" cy="1000000"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></am3d:spPr>
   <am3d:camera><am3d:pos x="0" y="0" z="1000000"/><am3d:up dx="0" dy="36000000" dz="0"/><am3d:lookAt x="0" y="0" z="0"/></am3d:camera>
   <am3d:raster><am3d:blip r:embed="rId3"/></am3d:raster>
  </am3d:model3d>
 </a:graphicData></a:graphic>
</p:graphicFrame>'''


def _write_native_3d_source(path: Path) -> Native3DAsset:
    poster = BytesIO()
    Image.new("RGB", (64, 64), (60, 140, 90)).save(poster, format="PNG")
    slide_xml = b'''<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cSld><p:spTree>''' + _MODEL_FRAME_XML + b'''</p:spTree></p:cSld></p:sld>'''
    rels = b'''<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
 <Relationship Id="rId2" Type="http://schemas.microsoft.com/office/2017/06/relationships/model3d" Target="../media/model3d1.glb"/>
 <Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/model-preview.png"/>
</Relationships>'''
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as package:
        package.writestr("ppt/slides/slide1.xml", slide_xml)
        package.writestr("ppt/slides/_rels/slide1.xml.rels", rels)
        package.writestr("ppt/media/model3d1.glb", b"glTF synthetic test payload")
        package.writestr("ppt/media/model-preview.png", poster.getvalue())
    return Native3DAsset(
        asset_ref="powerpoint-3d://fixture/1-test-cell",
        source_path=path,
        source_slide=1,
        name="3D Model 1",
        description="Test Cell",
        model_part="ppt/media/model3d1.glb",
        poster_part="ppt/media/model-preview.png",
        frame_xml=_MODEL_FRAME_XML,
    )


class PresentationMakerTests(unittest.TestCase):
    def setUp(self):
        maker._pending_presentation_parameters = None
        maker._pending_presentation_created_at = 0.0
        maker._presentation_quota_cooldown_until = 0.0
        presentation_assets._image_quota_cooldown_until = 0.0
        maker._presentation_output_locks.clear()

    def test_default_theme_is_restrained_mitsu_minimal(self):
        request = PresentationRequest.from_parameters({"topic": "Biology"})
        self.assertEqual(request.theme, "mitsu_minimal")
        self.assertEqual(request.resolved_appearance(), "dark")
        self.assertEqual(request.transition, "morph")

    def test_dark_and_light_appearance_follow_the_users_request(self):
        dark = PresentationRequest.from_parameters({
            "topic": "A dark cinematic presentation about astronomy",
        })
        light = PresentationRequest.from_parameters({
            "topic": "A bright presentation with a white background",
        })
        self.assertEqual(dark.resolved_appearance(), "dark")
        self.assertEqual(light.resolved_appearance(), "light")

    def test_request_infers_explicit_native_3d_model_intent(self):
        request = PresentationRequest.from_parameters({
            "topic": "Use the 3D models from this PowerPoint",
        })
        self.assertTrue(request.use_native_3d)

    def test_planner_fallback_marks_powerpoint_3d_requests(self):
        plan = planner._fallback_plan("Make a PowerPoint using its 3D models")
        parameters = plan["steps"][0]["parameters"]
        self.assertTrue(parameters["use_native_3d"])

    def test_native_powerpoint_3d_assets_are_discovered(self):
        with tempfile.TemporaryDirectory() as directory:
            source = Path(directory) / "source.pptx"
            _write_native_3d_source(source)
            assets = inspect_native_3d(source)

        self.assertEqual(len(assets), 1)
        self.assertEqual(assets[0].description, "Test Cell")
        self.assertEqual(assets[0].asset_ref, "powerpoint-3d://source/1-test-cell")

    def test_renderer_embeds_rotatable_native_powerpoint_3d(self):
        plan = maker._normalize_plan(_sample_plan(), "Cells", 6)
        with tempfile.TemporaryDirectory() as directory:
            source = Path(directory) / "source.pptx"
            asset = _write_native_3d_source(source)
            plan["slides"][0]["asset_ref"] = asset.asset_ref
            output = Path(directory) / "native-3d.pptx"
            maker._render_presentation(
                plan,
                output,
                native_3d_assets={asset.asset_ref: asset},
            )
            checks = verify_pptx(output, 6)

            self.assertEqual(checks["native_3d_model_count"], 1)
            self.assertEqual(checks["native_3d_shape_count"], 1)
            with zipfile.ZipFile(output) as package:
                content_types = package.read("[Content_Types].xml")
                slide_rels = package.read("ppt/slides/_rels/slide1.xml.rels")
            self.assertIn(b'model/gltf.binary', content_types)
            self.assertIn(b'relationships/model3d', slide_rels)

    def test_requested_native_3d_is_assigned_to_a_visual_slide(self):
        plan = maker._normalize_plan(_sample_plan(), "Cells", 6)
        with tempfile.TemporaryDirectory() as directory:
            source = Path(directory) / "source.pptx"
            asset = _write_native_3d_source(source)
            maker._enforce_native_3d_usage(
                plan,
                {asset.asset_ref: asset},
                requested=True,
            )
            placements = [
                slide for slide in plan["slides"]
                if slide.get("asset_ref") == asset.asset_ref
            ]
            output = Path(directory) / "reused-native-3d.pptx"
            maker._render_presentation(
                plan,
                output,
                native_3d_assets={asset.asset_ref: asset},
            )
            checks = verify_pptx(output, 6)

            self.assertGreaterEqual(len(placements), 3)
            self.assertTrue(all(slide["proof_type"] == "image" for slide in placements))
            self.assertGreaterEqual(checks["native_3d_shape_count"], 3)

    def test_render_creates_editable_widescreen_deck(self):
        with tempfile.TemporaryDirectory() as directory:
            output = Path(directory) / "clean-energy.pptx"
            maker._render_presentation(_sample_plan(), output, "arc_reactor", "Executives")

            self.assertTrue(output.exists())
            deck = Presentation(output)
            self.assertEqual(len(deck.slides), 6)
            self.assertAlmostEqual(deck.slide_width / deck.slide_height, 16 / 9, places=2)
            rendered_text = "\n".join(
                shape.text
                for slide in deck.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
            self.assertIn("The Future of Clean Energy", deck.core_properties.title)
            self.assertIn("Claim 3 moves the story forward", rendered_text)
            self.assertLessEqual(sum(len(slide.shapes) for slide in deck.slides), 34)

    def test_renderer_applies_morph_with_fade_fallback(self):
        with tempfile.TemporaryDirectory() as directory:
            output = Path(directory) / "morph.pptx"
            maker._render_presentation(_sample_plan(), output, transition="morph")
            checks = verify_pptx(output, 6)
            self.assertEqual(checks["transition_slide_count"], 5)
            self.assertEqual(checks["morph_slide_count"], 5)
            with zipfile.ZipFile(output) as package:
                second_slide = package.read("ppt/slides/slide2.xml")
            self.assertIn(b"morph", second_slide)
            self.assertIn(b"fade", second_slide)

    def test_morph_verification_survives_office_extension_parse_errors(self):
        with tempfile.TemporaryDirectory() as directory:
            output = Path(directory) / "morph-extension.pptx"
            maker._render_presentation(_sample_plan(), output, transition="morph")
            with patch(
                "actions.presentations.quality.ET.fromstring",
                side_effect=ET.ParseError("unknown Office extension"),
            ):
                checks = verify_pptx(output, 6)

        self.assertEqual(checks["transition_slide_count"], 5)
        self.assertEqual(checks["morph_slide_count"], 5)

    def test_missing_morph_marker_uses_fade_fallback_without_failing_deck(self):
        package = {
            "slide_count": 6,
            "chart_count": 0,
            "native_3d_model_count": 0,
            "native_3d_shape_count": 0,
            "native_3d_slide_count": 0,
            "transition_slide_count": 5,
            "morph_slide_count": 0,
            "empty_media": [],
            "bytes": 1000,
        }
        with (
            tempfile.TemporaryDirectory() as directory,
            patch.object(maker, "_plan_presentation", return_value=_sample_plan()),
            patch.object(maker, "verify_pptx", return_value=package),
        ):
            result = maker.build_presentation({
                "topic": "Cells",
                "slide_count": 6,
                "output_path": str(Path(directory) / "fade-fallback.pptx"),
                "quality": "fast",
                "export_pdf": False,
            })

        self.assertEqual(result.slide_count, 6)
        self.assertEqual(result.qa["transition_fallback"], "fade")
        self.assertTrue(any("Fade compatibility" in warning for warning in result.warnings))

    def test_same_output_path_uses_the_same_build_lock(self):
        path = Path("/tmp/mitsu-presentation-lock-test.pptx")
        self.assertIs(maker._output_lock_for(path), maker._output_lock_for(path))

    def test_renderer_uses_a_consistent_light_or_dark_canvas(self):
        with tempfile.TemporaryDirectory() as directory:
            light_path = Path(directory) / "light.pptx"
            dark_path = Path(directory) / "dark.pptx"
            maker._render_presentation(_sample_plan(), light_path, appearance="light")
            maker._render_presentation(_sample_plan(), dark_path, appearance="dark")
            light_fill = str(Presentation(light_path).slides[0].shapes[0].fill.fore_color.rgb)
            dark_fill = str(Presentation(dark_path).slides[0].shapes[0].fill.fore_color.rgb)
            self.assertEqual(light_fill, maker.THEMES["mitsu_minimal"]["paper"])
            self.assertEqual(dark_fill, maker.THEMES["mitsu_minimal"]["ink"])

    def test_model_focus_balances_native_3d_against_only_essential_text(self):
        plan = maker._normalize_plan(_sample_plan(), "Cells", 6)
        with tempfile.TemporaryDirectory() as directory:
            source = Path(directory) / "source.pptx"
            asset = _write_native_3d_source(source)
            maker._enforce_native_3d_usage(plan, {asset.asset_ref: asset}, requested=True)
            output = Path(directory) / "minimal-model-focus.pptx"
            maker._render_presentation(
                plan,
                output,
                native_3d_assets={asset.asset_ref: asset},
            )
            deck = Presentation(output)

            model_slides = []
            for slide in deck.slides:
                models = [shape for shape in slide.shapes if shape.name.startswith("!!MITSU 3D")]
                if models:
                    model_slides.append(slide)
                    posters = [shape for shape in slide.shapes if shape.name.startswith("!!MITSU 3D Poster")]
                    text_shapes = [
                        shape for shape in slide.shapes
                        if hasattr(shape, "text") and shape.text.strip()
                    ]
                    title = next(shape for shape in text_shapes if shape.name == "model-focus-title")
                    self.assertLess(models[0].left, Inches(1))
                    self.assertEqual(len(posters), 1)
                    self.assertGreater(title.left, Inches(7.5))
                    self.assertLessEqual(len(text_shapes), 4)
                    self.assertNotIn("First supporting idea", "\n".join(shape.text for shape in text_shapes))
                    self.assertEqual(title.text_frame.paragraphs[0].font.name, maker.DISPLAY_FONT)

        self.assertGreaterEqual(len(model_slides), 3)

    def test_normalizer_rejects_verbose_slide_copy(self):
        plan = _sample_plan()
        plan["slides"][1]["kicker"] = "A VERY LONG SECTION LABEL"
        plan["slides"][1]["title"] = "One two three four five six seven eight nine ten eleven"
        plan["slides"][1]["subtitle"] = " ".join(f"word{index}" for index in range(20))
        plan["slides"][1]["bullets"] = ["one", "two", "three"]
        normalized = maker._normalize_plan(plan, "Cells", 6)
        slide = normalized["slides"][1]

        self.assertLessEqual(len(slide["kicker"].split()), 2)
        self.assertLessEqual(len(slide["title"].split()), 8)
        self.assertLessEqual(len(slide["subtitle"].split()), 14)
        self.assertEqual(len(slide["bullets"]), 1)

    def test_every_slide_keeps_only_a_few_visible_text_regions(self):
        plan = maker._normalize_plan(_sample_plan(), "Minimal presentation", 6)
        with tempfile.TemporaryDirectory() as directory:
            output = Path(directory) / "minimal.pptx"
            maker._render_presentation(plan, output)
            deck = Presentation(output)
            counts = [
                len([
                    shape for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ])
                for slide in deck.slides
            ]
        self.assertTrue(all(count <= 4 for count in counts), counts)

    def test_feature_generates_and_stamps_output_without_network(self):
        with tempfile.TemporaryDirectory() as directory:
            output = Path(directory) / "generated.pptx"
            with patch.object(maker, "_plan_presentation", return_value=_sample_plan()):
                result = maker.create_presentation({
                    "topic": "clean energy",
                    "slide_count": 6,
                    "output_path": str(output),
                    "quality": "fast",
                    "export_pdf": False,
                })

            self.assertIn("Created an editable 6-slide", result)
            self.assertTrue(output.exists())
            self.assertTrue(output.with_name("generated.pptx.mitsu_meta.json").exists())

    def test_default_quality_caps_image_generation_at_two_requests(self):
        standard = PresentationRequest.from_parameters({"topic": "Biology"})
        premium = PresentationRequest.from_parameters({"topic": "Biology", "quality": "premium"})

        self.assertEqual(presentation_assets._image_request_limit(standard, 12), 2)
        self.assertEqual(presentation_assets._image_request_limit(premium, 12), 5)

    def test_image_generation_stops_immediately_when_quota_is_reached(self):
        plan = maker._normalize_plan(_sample_plan(), "Biology", 6)
        for slide in plan["slides"]:
            slide["type"] = "image" if slide["type"] not in {"cover", "closing"} else slide["type"]
            slide["image_prompt"] = "A clean scientific editorial image"
        request = PresentationRequest.from_parameters({"topic": "Biology"})
        client = types.SimpleNamespace(
            interactions=types.SimpleNamespace(
                create=Mock(side_effect=RuntimeError("429 RESOURCE_EXHAUSTED quota"))
            )
        )
        genai_module = types.ModuleType("google.genai")
        genai_module.Client = lambda **_kwargs: client
        google_module = types.ModuleType("google")
        google_module.genai = genai_module

        with tempfile.TemporaryDirectory() as directory, patch.dict(
            "sys.modules",
            {"google": google_module, "google.genai": genai_module},
        ):
            _lookup, warnings = presentation_assets.generate_planned_assets(
                plan,
                request,
                SourceBundle(),
                Path(directory),
                api_key="test-key",
            )

        self.assertEqual(client.interactions.create.call_count, 1)
        self.assertEqual(sum("quota" in warning.lower() for warning in warnings), 1)
        self.assertTrue(all(not slide.get("image_prompt") for slide in plan["slides"]))

        second_plan = maker._normalize_plan(_sample_plan(), "Biology", 6)
        for slide in second_plan["slides"]:
            slide["image_prompt"] = "Another scientific editorial image"
        with tempfile.TemporaryDirectory() as directory, patch.dict(
            "sys.modules",
            {"google": google_module, "google.genai": genai_module},
        ):
            presentation_assets.generate_planned_assets(
                second_plan,
                request,
                SourceBundle(),
                Path(directory),
                api_key="test-key",
            )
        self.assertEqual(client.interactions.create.call_count, 1)

    def test_planning_quota_falls_back_to_a_local_minimal_deck(self):
        with tempfile.TemporaryDirectory() as directory:
            output = Path(directory) / "quota-safe.pptx"
            second_output = Path(directory) / "quota-cooldown.pptx"
            with patch.object(
                maker,
                "_plan_presentation",
                side_effect=RuntimeError("429 RESOURCE_EXHAUSTED quota"),
            ) as planner:
                result = maker.build_presentation({
                    "topic": "Biology",
                    "slide_count": 6,
                    "output_path": str(output),
                    "quality": "quality",
                    "export_pdf": False,
                })
                second_result = maker.build_presentation({
                    "topic": "Ecology",
                    "slide_count": 6,
                    "output_path": str(second_output),
                    "quality": "quality",
                    "export_pdf": False,
                })

            deck = Presentation(output)
            text_counts = [
                sum(
                    1 for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                )
                for slide in deck.slides
            ]

        self.assertEqual(result.slide_count, 6)
        self.assertEqual(second_result.slide_count, 6)
        self.assertEqual(planner.call_count, 1)
        self.assertTrue(any("minimal local deck" in warning for warning in result.warnings))
        self.assertTrue(any("quota cooldown" in warning for warning in second_result.warnings))
        self.assertTrue(all(count <= 4 for count in text_counts), text_counts)

    def test_legacy_action_name_routes_to_create_presentation(self):
        with patch.object(maker, "create_presentation", return_value="created") as create:
            result = maker.presentation_maker({"topic": "clean energy"})

        self.assertEqual(result, "created")
        create.assert_called_once()

    def test_qa_output_paths_can_be_explicit(self):
        with tempfile.TemporaryDirectory() as directory:
            output = maker._output_path(str(Path(directory) / "briefing.pptx"), "Briefing")
            self.assertEqual(output, Path(directory) / "briefing.pptx")

    def test_filename_does_not_duplicate_powerpoint_extension(self):
        self.assertEqual(maker._safe_filename("Quarterly Review.pptx"), "Quarterly_Review.pptx")

    def test_live_tool_contract_declares_dedicated_action(self):
        project_root = Path(__file__).resolve().parent.parent
        tools = declared_tools(project_root / "main.py")
        self.assertEqual(tools.count("create_presentation"), 1)

    def test_planner_fallback_keeps_powerpoint_requests_on_dedicated_action(self):
        plan = planner._fallback_plan("Create a 12-slide PowerPoint about robotics")
        self.assertEqual(plan["steps"][0]["tool"], "create_presentation")
        self.assertEqual(plan["steps"][0]["parameters"]["slide_count"], 12)
        self.assertEqual(plan["steps"][0]["parameters"]["execution_mode"], "ask")

    def test_presentation_first_asks_about_native_3d(self):
        with patch.object(maker, "queue_presentation") as queued:
            message = maker.request_presentation({
                "topic": "A six-slide presentation about robotics",
                "slide_count": 6,
                "execution_mode": "ask",
            })

        self.assertEqual(message, maker.PRESENTATION_3D_QUESTION)
        self.assertFalse(queued.called)
        self.assertEqual(maker._pending_presentation_parameters["slide_count"], 6)

    def test_3d_answer_requires_an_actual_model_source(self):
        maker.request_presentation({
            "topic": "A presentation about robotics",
            "slide_count": 7,
            "execution_mode": "ask",
        })
        message = maker.request_presentation({
            "three_d_mode": "yes",
            "execution_mode": "ask",
        })
        self.assertEqual(message, maker.PRESENTATION_3D_SOURCE_QUESTION)
        self.assertTrue(maker._pending_presentation_parameters["use_native_3d"])
        self.assertEqual(maker._pending_presentation_parameters["slide_count"], 7)

    def test_attaching_a_model_deck_advances_to_the_run_mode_question(self):
        with tempfile.TemporaryDirectory() as directory:
            source = Path(directory) / "model-source.pptx"
            _write_native_3d_source(source)
            maker.request_presentation({
                "topic": "A presentation about robotics",
                "execution_mode": "ask",
            })
            maker.request_presentation({
                "three_d_mode": "yes",
                "execution_mode": "ask",
            })
            message = maker.request_presentation({
                "source_files": [str(source)],
                "execution_mode": "ask",
            })

        self.assertEqual(message, maker.PRESENTATION_RUN_MODE_QUESTION)
        self.assertEqual(
            maker._pending_presentation_parameters["model_source_file"],
            str(source),
        )
        self.assertNotIn("source_files", maker._pending_presentation_parameters)

    def test_model_library_does_not_turn_a_new_deck_into_a_redesign(self):
        request = PresentationRequest.from_parameters({
            "topic": "Robotics",
            "model_source_file": "/tmp/model-library.pptx",
        })

        self.assertEqual(request.resolved_mode(), "create")
        self.assertIsNone(request.source_deck())

    def test_direct_build_cannot_silently_drop_requested_3d(self):
        with tempfile.TemporaryDirectory() as directory:
            with (
                patch.object(maker, "_plan_presentation", return_value=_sample_plan()),
                self.assertRaisesRegex(ValueError, "contain no embedded 3D models"),
            ):
                maker.build_presentation({
                    "topic": "Robotics",
                    "use_native_3d": True,
                    "quality": "fast",
                    "export_pdf": False,
                    "output_path": str(Path(directory) / "must-not-build.pptx"),
                })

    def test_visible_choice_starts_the_pending_presentation_live(self):
        maker.request_presentation({
            "topic": "A presentation about robotics",
            "slide_count": 7,
            "execution_mode": "ask",
        })
        maker.request_presentation({"three_d_mode": "no", "execution_mode": "ask"})
        with patch.object(maker, "queue_presentation", return_value="live") as queued:
            message = maker.request_presentation({"execution_mode": "visible"})

        self.assertEqual(message, "live")
        parameters = queued.call_args.args[0]
        self.assertEqual(parameters["slide_count"], 7)
        self.assertEqual(parameters["topic"], "A presentation about robotics")
        self.assertTrue(queued.call_args.kwargs["visible"])

    def test_background_choice_starts_the_pending_presentation_quietly(self):
        maker.request_presentation({
            "topic": "A presentation about robotics",
            "execution_mode": "ask",
            "use_native_3d": False,
        })
        with patch.object(maker, "queue_presentation", return_value="background") as queued:
            message = maker.request_presentation({"execution_mode": "background"})

        self.assertEqual(message, "background")
        self.assertFalse(queued.call_args.kwargs["visible"])

    def test_visible_queue_publishes_real_presentation_phases(self):
        class Player:
            def __init__(self):
                self.shown = []
                self.updates = []

            def show_presentation_progress(self, title, visible=False):
                self.shown.append((title, visible))

            def update_presentation_progress(self, **update):
                self.updates.append(update)

            def write_log(self, _text):
                pass

        class Queue:
            def submit_job(self, **job):
                job["runner"](threading.Event(), lambda **_update: None)
                return "task-3d"

        class Result:
            artifacts = ["deck.pptx"]
            warnings = []

        player = Player()

        def build(*_args, **kwargs):
            kwargs["progress_callback"](percent=58, phase="Building editable PowerPoint")
            kwargs["progress_callback"](
                percent=100,
                phase="Presentation complete",
                artifacts=["deck.pptx"],
                warnings=[],
            )
            return Result()

        with (
            patch("agent.task_queue.get_queue", return_value=Queue()),
            patch.object(maker, "build_presentation", side_effect=build),
        ):
            message = maker.queue_presentation(
                {"topic": "Robotics", "export_pdf": False},
                player=player,
                visible=True,
            )

        self.assertIn("visible", message)
        self.assertEqual(player.shown, [("Robotics", True)])
        self.assertEqual(player.updates[0]["percent"], 58)
        self.assertTrue(player.updates[0]["visible"])
        self.assertEqual(player.updates[-1]["artifacts"], ["deck.pptx"])

    def test_request_accepts_multiple_sources_and_clamps_to_fifty_slides(self):
        request = PresentationRequest.from_parameters({
            "topic": "Operating review",
            "slide_count": 90,
            "source_file": "one.csv",
            "source_files": ["two.docx", "one.csv"],
            "quality": "premium",
        })
        self.assertEqual(request.slide_count, 50)
        self.assertEqual([path.name for path in request.source_files], ["two.docx", "one.csv"])
        self.assertEqual(request.quality, "premium")

    def test_csv_source_is_ingested_with_provenance(self):
        with tempfile.TemporaryDirectory() as directory:
            source = Path(directory) / "metrics.csv"
            source.write_text("quarter,revenue\nQ1,10\nQ2,12\n", encoding="utf-8")
            request = PresentationRequest.from_parameters({
                "topic": "Revenue",
                "source_files": [str(source)],
            })
            bundle = ingest_sources(request)
            self.assertIn("Q2, 12", bundle.text)
            self.assertEqual(bundle.records[0].provenance, "user-provided")

    def test_chart_slide_exports_native_editable_chart_part(self):
        plan = _sample_plan()
        plan["slides"][2]["type"] = "chart"
        plan["slides"][2]["chart"] = {
            "type": "column",
            "title": "Revenue",
            "categories": ["Q1", "Q2"],
            "series": [{"name": "Revenue", "values": [10, 12]}],
        }
        normalized = maker._normalize_plan(plan, "Revenue", 6)
        with tempfile.TemporaryDirectory() as directory:
            output = Path(directory) / "chart.pptx"
            maker._render_presentation(normalized, output)
            checks = verify_pptx(output, 6)
            self.assertEqual(checks["chart_count"], 1)
            with zipfile.ZipFile(output) as package:
                self.assertTrue(any(name.startswith("ppt/charts/chart") for name in package.namelist()))

    def test_image_led_plan_requires_real_visual_coverage(self):
        plan = maker._normalize_plan(_sample_plan(), "Cells", 6)
        plan["deck_profile"] = "consumer-retail"
        issues = maker._plan_quality_issues(plan)
        self.assertTrue(any("visual slides" in issue for issue in issues))
        self.assertTrue(any("cover requires" in issue for issue in issues))

    def test_visual_science_topic_is_routed_and_given_image_rhythm(self):
        plan = maker._normalize_plan(_sample_plan(), "Global warming and climate", 6)
        self.assertEqual(plan["deck_profile"], "consumer-retail")
        maker._enforce_visual_story(plan, "Global warming and climate")
        visual_slides = [
            slide for slide in plan["slides"]
            if slide.get("asset_ref") or slide.get("image_prompt")
        ]
        self.assertGreaterEqual(len(visual_slides), 3)
        self.assertTrue(plan["slides"][0]["image_prompt"])
        self.assertTrue(plan["slides"][-1]["image_prompt"])

    def test_cover_embeds_full_bleed_hero_asset(self):
        plan = maker._normalize_plan(_sample_plan(), "Cells", 6)
        plan["slides"][0]["asset_ref"] = "hero"
        with tempfile.TemporaryDirectory() as directory:
            hero = Path(directory) / "hero.png"
            Image.new("RGB", (1600, 900), (80, 160, 120)).save(hero)
            output = Path(directory) / "image-led.pptx"
            maker._render_presentation(plan, output, asset_lookup={"hero": hero})
            with zipfile.ZipFile(output) as package:
                media = [name for name in package.namelist() if name.startswith("ppt/media/")]
            self.assertTrue(media)

    def test_process_uses_one_clear_sequence_text_region(self):
        plan = maker._normalize_plan(_sample_plan(), "Energy", 6)
        with tempfile.TemporaryDirectory() as directory:
            output = Path(directory) / "process-safe.pptx"
            maker._render_presentation(plan, output)
            deck = Presentation(output)
            labels = sorted(
                (shape for shape in deck.slides[3].shapes if shape.name.startswith("process-label-")),
                key=lambda shape: shape.left,
            )
            self.assertEqual(len(labels), 1)

    def test_targeted_edit_never_overwrites_source_deck(self):
        with tempfile.TemporaryDirectory() as directory:
            source = Path(directory) / "source.pptx"
            output = Path(directory) / "edited.pptx"
            deck = Presentation()
            slide = deck.slides.add_slide(deck.slide_layouts[6])
            shape = slide.shapes.add_textbox(0, 0, 1000000, 500000)
            shape.text = "Old title"
            deck.save(source)
            edit_plan = {
                "title": "Edited",
                "deck_profile": "general",
                "slides": [],
                "edits": [{
                    "slide_number": 1,
                    "find_text": "Old title",
                    "replacement_text": "New title",
                }],
            }
            with patch.object(maker, "_plan_presentation", return_value=edit_plan):
                result = maker.build_presentation({
                    "topic": "Change Old title to New title",
                    "mode": "edit",
                    "source_file": str(source),
                    "output_path": str(output),
                    "quality": "fast",
                    "export_pdf": False,
                })
            source_text = "\n".join(shape.text for shape in Presentation(source).slides[0].shapes if hasattr(shape, "text"))
            output_text = "\n".join(shape.text for shape in Presentation(output).slides[0].shapes if hasattr(shape, "text"))
            self.assertIn("Old title", source_text)
            self.assertIn("New title", output_text)
            self.assertEqual(result.mode, "edit")

    def test_extend_preserves_source_slides_and_appends_to_requested_total(self):
        with tempfile.TemporaryDirectory() as directory:
            source = Path(directory) / "source.pptx"
            output = Path(directory) / "extended.pptx"
            deck = Presentation()
            for label in ("Existing one", "Existing two"):
                slide = deck.slides.add_slide(deck.slide_layouts[6])
                shape = slide.shapes.add_textbox(0, 0, 1000000, 500000)
                shape.text = label
            deck.save(source)
            append_plan = {
                "title": "Extended",
                "deck_profile": "general",
                "slides": [
                    {
                        "type": "statement",
                        "kicker": "NEXT",
                        "title": f"Added claim {index}",
                        "bullets": ["Evidence"],
                    }
                    for index in range(1, 4)
                ],
            }
            with patch.object(maker, "_plan_presentation", return_value=append_plan):
                result = maker.build_presentation({
                    "topic": "Extend the deck",
                    "mode": "extend",
                    "source_file": str(source),
                    "slide_count": 5,
                    "output_path": str(output),
                    "quality": "fast",
                    "export_pdf": False,
                })
            self.assertEqual(len(Presentation(source).slides), 2)
            self.assertEqual(len(Presentation(output).slides), 5)
            self.assertEqual(result.slide_count, 5)

    def test_specialized_queue_reports_progress_and_artifacts(self):
        queue = TaskQueue(max_concurrent=1)

        class Result:
            artifacts = ["presentation.pptx"]
            warnings = []

            def __str__(self):
                return "created"

        def runner(cancel_flag, progress):
            progress(percent=75, phase="Rendering")
            return Result()

        task_id = queue.submit_job("presentation", runner)
        task = queue._tasks[task_id]
        task.status = task.status.RUNNING
        queue._active_count = 1
        queue._run_task(task)
        status = queue.get_status(task_id)
        self.assertEqual(status["status"], "completed")
        self.assertEqual(status["progress"], 100)
        self.assertEqual(status["artifacts"], ["presentation.pptx"])

    def test_missing_pdf_renderer_delivers_powerpoint_with_warning(self):
        with tempfile.TemporaryDirectory() as directory:
            output = Path(directory) / "fallback.pptx"
            with (
                patch.object(maker, "_plan_presentation", return_value=_sample_plan()),
                patch.object(maker, "export_pdf", return_value=(None, "PDF unavailable")),
            ):
                result = maker.build_presentation({
                    "topic": "clean energy",
                    "slide_count": 6,
                    "output_path": str(output),
                    "quality": "fast",
                    "export_pdf": True,
                })
            self.assertTrue(output.exists())
            self.assertIsNone(result.pdf_path)
            self.assertIn("PDF unavailable", result.warnings)

    def test_cancelled_specialized_job_is_not_marked_failed(self):
        queue = TaskQueue(max_concurrent=1)

        def runner(cancel_flag, progress):
            cancel_flag.set()
            raise RuntimeError("cancelled")

        task_id = queue.submit_job("presentation", runner)
        task = queue._tasks[task_id]
        task.status = task.status.RUNNING
        queue._active_count = 1
        queue._run_task(task)
        status = queue.get_status(task_id)
        self.assertEqual(status["status"], "cancelled")
        self.assertEqual(status["phase"], "Cancelled")

    def test_build_honors_preexisting_cancel_flag_before_model_call(self):
        cancel = threading.Event()
        cancel.set()
        with tempfile.TemporaryDirectory() as directory:
            with self.assertRaisesRegex(RuntimeError, "cancelled"):
                maker.build_presentation({
                    "topic": "clean energy",
                    "output_path": str(Path(directory) / "cancelled.pptx"),
                    "quality": "fast",
                    "export_pdf": False,
                }, cancel_flag=cancel)

    def test_speaker_notes_are_editable_in_powerpoint(self):
        plan = maker._normalize_plan(_sample_plan(), "Energy", 6)
        plan["slides"][1]["speaker_notes"] = "Explain the transition assumptions."
        with tempfile.TemporaryDirectory() as directory:
            output = Path(directory) / "notes.pptx"
            maker._render_presentation(plan, output, include_speaker_notes=True)
            deck = Presentation(output)
            self.assertIn("transition assumptions", deck.slides[1].notes_slide.notes_text_frame.text)


if __name__ == "__main__":
    unittest.main()
