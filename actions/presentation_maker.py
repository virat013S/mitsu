"""Create polished, editable PowerPoint presentations for MITSU users."""

from __future__ import annotations

import json
import mimetypes
import os
import platform
import re
import shutil
import subprocess
import threading
import time
import uuid
from pathlib import Path
from typing import Any, Callable

from actions.mitsu_file_stamp import generated_by, mark_created_file
from actions.presentations.assets import generate_planned_assets
from actions.presentations.models import (
    MAX_SLIDES,
    MIN_SLIDES,
    PresentationRequest,
    PresentationResult,
    SourceBundle,
)
from actions.presentations.models3d import (
    Native3DAsset,
    add_native_3d_model,
    inspect_native_3d,
    native_3d_lookup,
)
from actions.presentations.quality import (
    clean_preview_dir,
    create_contact_sheet,
    critique_renders,
    export_pdf,
    render_pdf,
    verify_pptx,
)
from actions.presentations.sources import ingest_sources, source_manifest


BASE_DIR = Path(__file__).resolve().parent.parent
API_CONFIG_PATH = BASE_DIR / "config" / "api_keys.json"
DEFAULT_MODEL = "gemini-2.5-flash"
PRESENTATION_PENDING_TTL_SECONDS = 15 * 60
PRESENTATION_QUOTA_COOLDOWN_SECONDS = 5 * 60
PRESENTATION_RUN_MODE_QUESTION = (
    "I have the presentation brief ready, but I have not started building it. "
    "Would you like to see the task while I work, or keep it running in the background?"
)
PRESENTATION_3D_QUESTION = (
    "Before I build it, would you like a native, rotatable 3D model in the presentation? "
    "If you say yes, I will use a relevant model embedded in the PowerPoint sources you provide."
)
PRESENTATION_3D_SOURCE_QUESTION = (
    "Yes. I need the actual model before I start. Please attach or select a PowerPoint file "
    "that contains a relevant 3D model. You can also say continue without 3D."
)
_pending_presentation_lock = threading.Lock()
_pending_presentation_parameters: dict[str, Any] | None = None
_pending_presentation_created_at = 0.0
_tenant_pending_presentations: dict[str, tuple[dict[str, Any], float]] = {}
_presentation_quota_cooldown_until = 0.0
_presentation_output_locks_guard = threading.Lock()
_presentation_output_locks: dict[str, threading.Lock] = {}

THEMES = {
    "mitsu_minimal": {
        "ink": "020A12", "paper": "F3F7F8", "accent": "19DDF4",
        "accent_2": "88F4FF", "muted": "788793", "white": "EAFBFF",
        "trace": "102B36",
    },
    "editorial": {
        "ink": "10171B", "paper": "F4F0E8", "accent": "76C893",
        "accent_2": "E07A5F", "muted": "68777D", "white": "FFFFFF",
    },
    "arc_reactor": {
        "ink": "07141F", "paper": "F3EFE6", "accent": "39D9FF",
        "accent_2": "FFB547", "muted": "8193A3", "white": "FFFFFF",
    },
    "executive": {
        "ink": "111722", "paper": "F5F1E8", "accent": "D4553D",
        "accent_2": "C5A46D", "muted": "727985", "white": "FFFFFF",
    },
    "platinum": {
        "ink": "17191D", "paper": "F4F5F6", "accent": "6D7C91",
        "accent_2": "B5854B", "muted": "767D87", "white": "FFFFFF",
    },
}

# The MITSU application bundles Space Grotesk and JetBrains Mono for Qt, but
# PowerPoint cannot see application-private fonts. These native macOS families
# preserve the same restrained geometric/technical hierarchy without fallback.
DISPLAY_FONT = "Avenir Next"
BODY_FONT = "Avenir Next"
DATA_FONT = "Menlo"

DECK_PROFILES = {
    "finance-ir", "product-platform", "gtm-growth", "engineering-platform",
    "strategy-leadership", "consumer-retail", "appendix-heavy", "general",
}

_IMAGE_LED_TOPIC_RE = re.compile(
    r"\b(?:biology|cell|cells|animal|animals|plant|plants|nature|climate|global warming|"
    r"pollution|environment|ecosystem|ocean|space|planet|travel|destination|food|fashion|"
    r"beauty|sport|architecture|interior|art|museum|history|geography)\b",
    re.IGNORECASE,
)


def _api_key() -> str:
    from memory.config_manager import get_gemini_key

    key = str(get_gemini_key() or "").strip()
    if not key:
        raise RuntimeError("GEMINI_API_KEY is not configured.")
    return key


def _is_quota_error(error: Exception) -> bool:
    message = str(error).lower()
    return any(token in message for token in ("429", "quota", "resource_exhausted", "rate limit"))


def _presentation_quota_cooldown_active() -> bool:
    return time.monotonic() < _presentation_quota_cooldown_until


def _mark_presentation_quota_cooldown() -> None:
    global _presentation_quota_cooldown_until
    _presentation_quota_cooldown_until = (
        time.monotonic() + PRESENTATION_QUOTA_COOLDOWN_SECONDS
    )


def _output_lock_for(path: Path) -> threading.Lock:
    key = str(path.expanduser().resolve())
    with _presentation_output_locks_guard:
        return _presentation_output_locks.setdefault(key, threading.Lock())


def _safe_filename(value: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._ -]+", "", str(value or "")).strip(" .")
    cleaned = re.sub(r"\s+", "_", cleaned)
    cleaned = cleaned[:80] or "MITSU_Presentation"
    return cleaned if cleaned.lower().endswith(".pptx") else cleaned + ".pptx"


def _output_path(raw_path: str, title: str) -> Path:
    raw = str(raw_path or "").strip().strip('"').strip("'")
    if not raw:
        return Path.home() / "Desktop" / "MITSU Presentations" / _safe_filename(title)
    path = Path(raw).expanduser()
    if not path.is_absolute():
        path = Path.home() / "Desktop" / path
    if path.suffix.lower() != ".pptx":
        path = path / _safe_filename(title) if not path.suffix else path.with_suffix(".pptx")
    return path


def _extract_source(path_value: str) -> str:
    raw = str(path_value or "").strip().strip('"').strip("'")
    if not raw:
        return ""
    path = Path(raw).expanduser()
    if not path.exists() or not path.is_file():
        raise FileNotFoundError(f"Presentation source file was not found: {path}")

    suffix = path.suffix.lower()
    if suffix in {".txt", ".md", ".rst", ".csv", ".json", ".jsonl"}:
        text = path.read_text(encoding="utf-8", errors="replace")
    elif suffix == ".docx":
        from docx import Document

        text = "\n".join(p.text for p in Document(path).paragraphs if p.text.strip())
    elif suffix == ".pptx":
        from pptx import Presentation

        chunks = []
        for slide in Presentation(path).slides:
            chunks.extend(
                shape.text for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            )
        text = "\n".join(chunks)
    else:
        raise ValueError("Source files must be TXT, Markdown, CSV, JSON, DOCX, or PPTX.")
    return text[:50_000]


def _clean_json(text: str) -> dict[str, Any]:
    value = str(text or "").strip()
    value = re.sub(r"^```(?:json)?\s*", "", value, flags=re.IGNORECASE)
    value = re.sub(r"\s*```$", "", value)
    start, end = value.find("{"), value.rfind("}")
    if start < 0 or end <= start:
        raise ValueError("Gemini did not return a presentation plan.")
    result = json.loads(value[start:end + 1])
    if not isinstance(result, dict):
        raise ValueError("Presentation plan was not a JSON object.")
    return result


def _plan_presentation(
    topic: str,
    title: str,
    audience: str,
    slide_count: int,
    tone: str,
    source_text: str,
    *,
    request: PresentationRequest | None = None,
    bundle: SourceBundle | None = None,
    cancel_flag: Any = None,
) -> dict[str, Any]:
    from google import genai
    from google.genai import types

    mode = request.resolved_mode() if request else "create"
    quality = request.quality if request else "quality"
    language = request.language if request else ""
    include_notes = bool(request and request.include_speaker_notes)
    appearance = request.resolved_appearance() if request else "light"
    source_block = (
        "Use the source material below as the factual authority. Do not invent "
        "statistics, quotations, customers, or citations.\n\nSOURCE MATERIAL:\n" + source_text
        if source_text else
        "No source material was supplied. Use stable general knowledge and avoid precise, unsupported statistics."
    )
    mode_instruction = ""
    if mode == "edit":
        mode_instruction = """
This is a targeted edit of an existing PowerPoint. Return an `edits` array with
operation (replace_text | style_text | delete_slide), slide_number (1-based),
find_text, replacement_text, optional font_color (six-digit hex), font_size,
bold, and notes. Preserve all unrelated slides, masters, layouts, media, and
animations. Return `slides` as an empty array unless the instruction explicitly
asks to add slides.
"""
    elif mode == "extend":
        mode_instruction = """
This extends an existing PowerPoint. Preserve the source deck and return only the
new slides to append. They must inherit the source deck's visual and narrative grammar.
"""
    elif mode == "redesign":
        mode_instruction = """
This redesigns an existing deck. Use its facts and brand cues but rebuild the
narrative and slide compositions. Do not fabricate or approximate identity assets.
"""

    prompt = f"""You are MITSU, an elite presentation editor and data storyteller.
Task mode: {mode}
Quality mode: {quality}
Presentation appearance: {appearance}
Default slide transition: {request.transition if request else 'morph'}
Create a coherent {slide_count}-slide presentation plan about: {topic}
Working title: {title or topic}
Audience: {audience or 'general professional audience'}
Tone: {tone or 'confident, concise, executive'}
Output language: {language or 'match the user request and sources'}
{mode_instruction}

Choose exactly one deck_profile: finance-ir | product-platform | gtm-growth |
engineering-platform | strategy-leadership | consumer-retail | appendix-heavy | general.
Every non-appendix slide must make one specific claim and have one dominant proof
object. Use at least five macro-layout families in a 10-slide deck, never repeat
the same layout three times, and avoid generic card grids. Never invent metrics.
Every slide must be minimal and immediately readable: one idea, one dominant
visual, and normally only three visible text regions: a tiny section marker, a
short claim title, and one essential support line. Every visible word must add
topic-specific information. Omit optional copy instead of filling empty space.
Never write generic filler such as "overview", "key points", "what matters",
"this section", or "next steps". Do not scatter numbered fragments or small text
boxes across the canvas. Use at most one concise bullet unless a data table
genuinely requires more. The visual and text must balance as one composition;
clean negative space is intentional and must not be filled with decoration.
Use the consumer-retail profile for educational science, nature, travel, product,
people, place, animal, food, art, and other visually inspectable subjects. Those
decks must be image-led: give the cover an image_prompt and give at least half of
all slides a unique asset_ref or image_prompt. Describe a different crop, subject,
and visual role for every generated image. Never use an empty visual placeholder.
Prefer full-bleed detail crops, editorial split compositions, and section-scale
typography. Body copy must remain readable over imagery; protect it with a solid
contrast field. Do not imitate tiny low-contrast text from a reference deck.
Follow the reference deck's dimensional editorial idea: large visual subjects,
purposeful negative space, changing scale and crop, varied compositions, and a
clear visual journey instead of repeated panels. Honor the requested light or
dark appearance consistently across the deck.
When native PowerPoint 3D is requested, make the model the dominant object on
roughly one half of the canvas and place the short text block in the open half.
Do not put text over the model, wrap it in a card, or add extra explanatory copy.
For create/redesign, slide 1 must be "cover" and the final slide "closing".
Return exactly {slide_count} slides unless task mode is edit. Each slide may contain:
- type: cover | statement | two_column | process | comparison | chart | table |
  timeline | image | closing | appendix
- kicker: 1-2 uppercase words
- title: a specific claim, maximum 8 words
- subtitle: one essential sentence, maximum 14 words; blank when unnecessary
- bullets: 0-1 concise strings
- proof_type: chart | table | process | comparison | timeline | image | statement
- visual_label and visual_items: concise diagram or sequence content
- left_label and right_label: concise comparison labels when relevant
- left_items and right_items: up to 4 strings each when relevant
- chart: {{"type":"bar|column|line|area|pie|doughnut|scatter", "title":"...",
  "categories":["..."], "series":[{{"name":"...","values":[1,2]}}]}}
- table: {{"headers":["..."], "rows":[["..."]]}}
- asset_ref: exact supplied image filename or exact `powerpoint-3d://...` asset ID
  when the slide should use supplied media. A PowerPoint 3D asset stays editable
  and rotatable; use it when the user explicitly asks for 3D and the source list
  contains a semantically relevant native model. Never invent a 3D asset ID.
- image_prompt: an editorial image prompt only when a generated image materially helps;
  never request text, logos, product UI, mascots, badges, or brand marks
- layout_variant: auto | full_bleed | image_left | image_right | editorial
- speaker_notes: useful presenter notes when requested
- source_note: identify the supporting source number or URL; blank if unsupported

For edit mode also return edits. For every mode return citations as an array and
design_system with palette intent, typography intent, chart grammar, and banned motifs.

Return JSON only in this shape:
{{"title":"...","subtitle":"...","deck_profile":"general",
"design_system":{{}},"citations":[],"edits":[],"slides":[{{...}}]}}

{source_block}
"""
    client = genai.Client(api_key=_api_key())
    contents: list[Any] = [prompt]
    uploaded_names: list[str] = []
    if bundle:
        for media_path in bundle.media_paths[:10]:
            _check_cancel(cancel_flag)
            mime_type = mimetypes.guess_type(media_path)[0] or "application/octet-stream"
            try:
                if media_path.stat().st_size <= 18_000_000:
                    contents.append(types.Part.from_bytes(data=media_path.read_bytes(), mime_type=mime_type))
                else:
                    uploaded = client.files.upload(file=str(media_path))
                    if getattr(uploaded, "name", None):
                        uploaded_names.append(uploaded.name)
                    deadline = time.time() + 180
                    while getattr(getattr(uploaded, "state", None), "name", "ACTIVE") == "PROCESSING":
                        _check_cancel(cancel_flag)
                        if time.time() >= deadline:
                            raise TimeoutError(f"Timed out processing {media_path.name}")
                        time.sleep(2)
                        uploaded = client.files.get(name=uploaded.name)
                    contents.append(uploaded)
            except Exception as exc:
                bundle.warnings.append(f"Could not analyze media source {media_path.name}: {exc}")
    try:
        response = client.models.generate_content(
            model=DEFAULT_MODEL,
            contents=contents,
            config={"response_mime_type": "application/json"},
        )
        return _clean_json(response.text)
    finally:
        for name in uploaded_names:
            try:
                client.files.delete(name=name)
            except Exception:
                pass


def _trim(value: Any, maximum: int) -> str:
    text = " ".join(str(value or "").split())
    return text if len(text) <= maximum else text[:maximum - 1].rstrip() + "…"


def _concise(value: Any, maximum_words: int, maximum_chars: int) -> str:
    text = _trim(value, maximum_chars)
    words = text.split()
    if len(words) <= maximum_words:
        return text
    return " ".join(words[:maximum_words]).rstrip(".,;:") + "…"


def _items(value: Any, limit: int = 5, maximum: int = 90) -> list[str]:
    if not isinstance(value, list):
        return []
    return [_trim(item, maximum) for item in value[:limit] if str(item or "").strip()]


def _chart(value: Any) -> dict[str, Any]:
    raw = value if isinstance(value, dict) else {}
    chart_type = str(raw.get("type") or "column").lower()
    allowed = {"bar", "column", "line", "area", "pie", "doughnut", "scatter"}
    if chart_type not in allowed:
        chart_type = "column"
    categories = _items(raw.get("categories"), 15, 32)
    series = []
    for item in raw.get("series", []) if isinstance(raw.get("series"), list) else []:
        if not isinstance(item, dict):
            continue
        values = []
        for number in item.get("values", [])[: len(categories) or 15]:
            try:
                values.append(float(number))
            except (TypeError, ValueError):
                values.append(0.0)
        if values:
            series.append({"name": _trim(item.get("name") or "Series", 40), "values": values})
    return {
        "type": chart_type,
        "title": _trim(raw.get("title"), 80),
        "categories": categories,
        "series": series[:6],
    }


def _table(value: Any) -> dict[str, Any]:
    raw = value if isinstance(value, dict) else {}
    headers = _items(raw.get("headers"), 8, 35)
    rows = []
    raw_rows = raw.get("rows") if isinstance(raw.get("rows"), list) else []
    for row in raw_rows[:12]:
        if isinstance(row, list):
            rows.append([_trim(cell, 45) for cell in row[: len(headers) or 8]])
    return {"headers": headers, "rows": rows}


def _normalize_plan(
    plan: dict[str, Any],
    topic: str,
    count: int,
    *,
    bookends: bool = True,
) -> dict[str, Any]:
    raw_slides = plan.get("slides") if isinstance(plan.get("slides"), list) else []
    if len(raw_slides) != count:
        raise ValueError(f"Gemini returned {len(raw_slides)} slides; expected {count}.")
    allowed = {
        "cover", "statement", "two_column", "process", "comparison", "chart",
        "table", "timeline", "image", "closing", "appendix",
    }
    slides = []
    for index, raw in enumerate(raw_slides):
        raw = raw if isinstance(raw, dict) else {}
        slide_type = str(raw.get("type") or "statement").lower()
        if bookends and index == 0:
            slide_type = "cover"
        elif bookends and index == count - 1:
            slide_type = "closing"
        elif slide_type not in allowed or slide_type in {"cover", "closing"}:
            slide_type = (
                "statement", "two_column", "chart", "process", "comparison", "table", "timeline",
            )[(index - 1) % 7]
        slides.append({
            "type": slide_type,
            "kicker": _concise(raw.get("kicker") or f"SECTION {index}", 2, 20).upper(),
            "title": _concise(raw.get("title") or topic, 8, 72),
            "subtitle": _concise(raw.get("subtitle"), 14, 110),
            "bullets": _items(raw.get("bullets"), 1, 80),
            "visual_label": _trim(raw.get("visual_label") or "KEY IDEA", 45).upper(),
            "visual_items": _items(raw.get("visual_items"), 3, 55),
            "left_label": _trim(raw.get("left_label") or "TODAY", 30).upper(),
            "right_label": _trim(raw.get("right_label") or "NEXT", 30).upper(),
            "left_items": _items(raw.get("left_items"), 2, 65),
            "right_items": _items(raw.get("right_items"), 2, 65),
            "proof_type": _trim(raw.get("proof_type") or slide_type, 25).lower(),
            "chart": _chart(raw.get("chart")),
            "table": _table(raw.get("table")),
            "asset_ref": _trim(raw.get("asset_ref"), 180),
            "image_prompt": _trim(raw.get("image_prompt"), 500),
            "layout_variant": (
                str(raw.get("layout_variant") or "auto").lower()
                if str(raw.get("layout_variant") or "auto").lower()
                in {"auto", "full_bleed", "image_left", "image_right", "editorial"}
                else "auto"
            ),
            "speaker_notes": _trim(raw.get("speaker_notes"), 1500),
            "source_note": _trim(raw.get("source_note"), 140),
        })
    profile = str(plan.get("deck_profile") or "general").lower()
    if profile not in DECK_PROFILES:
        profile = "general"
    if profile in {"general", "strategy-leadership"} and _IMAGE_LED_TOPIC_RE.search(topic):
        profile = "consumer-retail"
    return {
        "title": _trim(plan.get("title") or topic, 100),
        "subtitle": _trim(plan.get("subtitle"), 180),
        "deck_profile": profile,
        "design_system": plan.get("design_system") if isinstance(plan.get("design_system"), dict) else {},
        "citations": _items(plan.get("citations"), 30, 220),
        "edits": plan.get("edits") if isinstance(plan.get("edits"), list) else [],
        "slides": slides,
    }


def _quota_fallback_plan(
    request: PresentationRequest,
    count: int,
    source_text: str,
) -> dict[str, Any]:
    """Create a restrained, editable deck when planning quota is unavailable."""
    topic = _trim(request.title or request.topic, 90)
    evidence: list[str] = []
    for raw in str(source_text or "").splitlines():
        line = " ".join(raw.split()).strip(" -\t")
        if (
            20 <= len(line) <= 150
            and not re.fullmatch(r"(?:SLIDE|SHEET)\s+\d+", line, re.IGNORECASE)
            and line.casefold() not in {item.casefold() for item in evidence}
        ):
            evidence.append(line)
        if len(evidence) >= 10:
            break

    def evidence_at(index: int, fallback: str) -> str:
        return evidence[index % len(evidence)] if evidence else fallback

    middle_types = ("statement", "two_column", "process", "comparison")
    slides: list[dict[str, Any]] = []
    for index in range(count):
        if index == 0:
            slides.append({
                "type": "cover",
                "kicker": "MITSU BRIEFING",
                "title": topic,
                "subtitle": evidence_at(0, "A clear, focused presentation built around the essential ideas."),
            })
            continue
        if index == count - 1:
            slides.append({
                "type": "closing",
                "kicker": "CONCLUSION",
                "title": "Keep the conclusion clear and actionable",
                "subtitle": evidence_at(index, f"Use the evidence to decide what matters next for {topic}."),
            })
            continue

        slide_type = middle_types[(index - 1) % len(middle_types)]
        support = evidence_at(index, f"This section keeps the discussion of {topic} focused and easy to follow.")
        slide = {
            "type": slide_type,
            "kicker": ("CONTEXT", "FOCUS", "SEQUENCE", "CHOICES")[(index - 1) % 4],
            "title": (
                "Start with the question that matters",
                "Separate the essential ideas from the noise",
                "Move from understanding to practical action",
                "Make the important trade-offs visible",
            )[(index - 1) % 4],
            "subtitle": support,
            "bullets": [support],
            "visual_label": "KEY STEPS",
            "visual_items": ["Frame the question", "Review the evidence", "Choose the next step"],
            "left_label": "CONTEXT",
            "right_label": "DIRECTION",
            "left_items": ["What is known", "What remains unclear"],
            "right_items": ["What matters now", "What should happen next"],
        }
        slides.append(slide)
    return {
        "title": topic,
        "subtitle": "A concise MITSU briefing",
        "deck_profile": "general",
        "design_system": {},
        "citations": [],
        "edits": [],
        "slides": slides,
    }


def _enforce_visual_story(plan: dict[str, Any], topic: str) -> dict[str, Any]:
    """Guarantee an image-led rhythm when the subject is visually inspectable."""
    if plan.get("deck_profile") != "consumer-retail":
        return plan
    slides = plan.get("slides") if isinstance(plan.get("slides"), list) else []
    if not slides:
        return plan
    target = max(2, (len(slides) + 1) // 2)
    variants = ("image_right", "full_bleed", "image_left", "editorial")

    def ensure_prompt(slide: dict, role: str, index: int) -> None:
        if slide.get("asset_ref") or slide.get("image_prompt"):
            return
        title = str(slide.get("title") or topic).strip()
        slide["image_prompt"] = (
            f"Editorial widescreen visual for {topic}. Slide role: {role}. "
            f"Show the concrete subject behind this claim: {title}. "
            "One clear focal subject, dimensional lighting, strong crop, generous text-safe negative space, "
            "no words, labels, logos, watermarks, UI, or infographic text."
        )
        slide["layout_variant"] = variants[index % len(variants)]

    ensure_prompt(slides[0], "opening hero", 0)
    if len(slides) > 1:
        ensure_prompt(slides[-1], "closing synthesis", len(slides) - 1)

    visual_count = sum(bool(slide.get("asset_ref") or slide.get("image_prompt")) for slide in slides)
    candidates = [
        (index, slide) for index, slide in enumerate(slides[1:-1], start=1)
        if slide.get("type") in {"image", "statement", "two_column"}
    ]
    for index, slide in candidates:
        if visual_count >= target:
            break
        if slide.get("asset_ref") or slide.get("image_prompt"):
            continue
        slide["type"] = "image"
        ensure_prompt(slide, "visual evidence", index)
        visual_count += 1
    return plan


def _enforce_native_3d_usage(
    plan: dict[str, Any],
    assets: dict[str, Native3DAsset],
    requested: bool,
) -> dict[str, Any]:
    """Carry real native models through a substantial part of the visual story."""
    if not requested or not assets:
        return plan
    slides = plan.get("slides") if isinstance(plan.get("slides"), list) else []
    if not slides:
        return plan
    native_refs = {
        asset.asset_ref.lower(): asset
        for asset in assets.values()
    }
    unique_assets = list(native_refs.values())
    existing_indices = {
        index for index, slide in enumerate(slides)
        if str(slide.get("asset_ref") or "").lower() in native_refs
    }
    image_candidates = [
        (index, slide) for index, slide in enumerate(slides)
        if slide.get("type") == "image" and index not in existing_indices
    ]
    narrative_candidates = [
        (index, slide) for index, slide in enumerate(slides[1:-1], start=1)
        if slide.get("type") in {"statement", "two_column", "process"}
        and index not in existing_indices
        and all(index != candidate_index for candidate_index, _ in image_candidates)
    ]
    remaining_candidates = [
        (index, slide) for index, slide in enumerate(slides)
        if index not in existing_indices
        and all(
            index != candidate_index
            for candidate_index, _ in image_candidates + narrative_candidates
        )
    ]
    candidates = image_candidates + narrative_candidates + remaining_candidates

    target = min(len(slides), max(2, (len(slides) + 1) // 2))
    needed = max(0, target - len(existing_indices))
    variants = ("image_right", "image_left", "full_bleed")
    for placement, (index, slide) in enumerate(candidates[:needed]):
        asset = unique_assets[(len(existing_indices) + placement) % len(unique_assets)]
        if slide.get("type") not in {"cover", "closing"}:
            slide["type"] = "image"
        slide["asset_ref"] = asset.asset_ref
        slide["image_prompt"] = ""
        slide["layout_variant"] = variants[placement % len(variants)]
        slide["proof_type"] = "image"
    return plan


def _request_native_3d_assets(parameters: dict) -> list[Native3DAsset]:
    raw_paths: list[Any] = []
    for key in ("source_file", "source_files", "template_file", "model_source_file"):
        value = parameters.get(key)
        if isinstance(value, (list, tuple, set)):
            raw_paths.extend(value)
        elif value:
            raw_paths.append(value)
    assets: list[Native3DAsset] = []
    seen_paths: set[str] = set()
    for raw in raw_paths:
        path = Path(str(raw)).expanduser()
        key = str(path.resolve()) if path.exists() else str(path)
        if key in seen_paths or path.suffix.lower() != ".pptx":
            continue
        seen_paths.add(key)
        assets.extend(inspect_native_3d(path))
    return assets


def _apply_slide_transition(slide, transition: str = "morph") -> None:
    """Apply a native PowerPoint transition, preferring Morph with Fade fallback."""
    from pptx.oxml import parse_xml

    value = str(transition or "morph").strip().lower()
    root = slide._element
    for child in list(root):
        local_name = str(child.tag).rsplit("}", 1)[-1]
        if local_name == "transition":
            root.remove(child)
        elif local_name == "AlternateContent" and any(
            str(node.tag).rsplit("}", 1)[-1] == "transition"
            for node in child.iter()
        ):
            root.remove(child)
    if value == "none":
        return
    if value == "fade":
        root.append(parse_xml(
            '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            'spd="med"><p:fade/></p:transition>'
        ))
        return
    root.append(parse_xml(
        '<mc:AlternateContent '
        'xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" '
        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" '
        'xmlns:p159="http://schemas.microsoft.com/office/powerpoint/2015/09/main">'
        '<mc:Choice Requires="p159">'
        '<p:transition spd="slow" p14:dur="1500">'
        '<p159:morph option="byObject"/>'
        '</p:transition>'
        '</mc:Choice>'
        '<mc:Fallback>'
        '<p:transition spd="med"><p:fade/></p:transition>'
        '</mc:Fallback>'
        '</mc:AlternateContent>'
    ))


def _normalize_edit_plan(plan: dict[str, Any], topic: str) -> dict[str, Any]:
    edits = []
    for raw in plan.get("edits", []) if isinstance(plan.get("edits"), list) else []:
        if not isinstance(raw, dict):
            continue
        try:
            slide_number = int(raw.get("slide_number"))
        except (TypeError, ValueError):
            continue
        find_text = _trim(raw.get("find_text"), 500)
        replacement = _trim(raw.get("replacement_text"), 2000)
        operation = str(raw.get("operation") or "replace_text").lower()
        if operation not in {"replace_text", "style_text", "delete_slide"}:
            operation = "replace_text"
        if slide_number > 0 and (find_text or operation == "delete_slide"):
            edits.append({
                "operation": operation,
                "slide_number": slide_number,
                "find_text": find_text,
                "replacement_text": replacement,
                "font_color": re.sub(r"[^0-9A-Fa-f]", "", str(raw.get("font_color") or ""))[:6],
                "font_size": raw.get("font_size"),
                "bold": raw.get("bold"),
                "notes": _trim(raw.get("notes"), 500),
            })
    result = {
        "title": _trim(plan.get("title") or topic, 100),
        "subtitle": _trim(plan.get("subtitle"), 180),
        "deck_profile": str(plan.get("deck_profile") or "general").lower(),
        "design_system": plan.get("design_system") if isinstance(plan.get("design_system"), dict) else {},
        "citations": _items(plan.get("citations"), 30, 220),
        "edits": edits,
        "slides": [],
    }
    raw_slides = plan.get("slides") if isinstance(plan.get("slides"), list) else []
    if raw_slides:
        synthetic = dict(plan)
        synthetic["slides"] = raw_slides
        result["slides"] = _normalize_plan(
            synthetic, topic, len(raw_slides), bookends=False,
        )["slides"]
    return result


def _plan_quality_issues(plan: dict[str, Any]) -> list[str]:
    """Catch weak narrative plans before spending time rendering them."""
    slides = plan.get("slides") if isinstance(plan.get("slides"), list) else []
    issues: list[str] = []
    titles: set[str] = set()
    repeated = 1
    previous_type = ""
    seen_types: set[str] = set()

    for index, slide in enumerate(slides, start=1):
        slide_type = str(slide.get("type") or "")
        title = str(slide.get("title") or "").strip()
        seen_types.add(slide_type)
        if title.lower() in titles:
            issues.append(f"Slide {index} repeats an earlier title.")
        titles.add(title.lower())
        if len(title.split()) > 14:
            issues.append(f"Slide {index} title is too long to scan.")

        repeated = repeated + 1 if slide_type == previous_type else 1
        previous_type = slide_type
        if repeated >= 3:
            issues.append(f"Slide {index} repeats the same layout three times in a row.")

        if slide_type in {"cover", "closing", "appendix"}:
            continue
        has_proof = {
            "chart": bool(slide.get("chart", {}).get("categories") and slide.get("chart", {}).get("series")),
            "table": bool(slide.get("table", {}).get("headers") and slide.get("table", {}).get("rows")),
            "process": bool(slide.get("visual_items")),
            "timeline": bool(slide.get("visual_items")),
            "comparison": bool(slide.get("left_items") and slide.get("right_items")),
            "image": bool(slide.get("asset_ref") or slide.get("image_prompt")),
            "two_column": bool(slide.get("bullets") and slide.get("visual_items")),
            "statement": bool(slide.get("bullets") or slide.get("visual_items") or slide.get("subtitle")),
        }.get(slide_type, bool(slide.get("bullets") or slide.get("visual_items")))
        if not has_proof:
            issues.append(f"Slide {index} makes a claim without a dominant proof object.")

    if len(slides) >= 8 and len(seen_types) < 4:
        issues.append("The deck needs more layout variety.")
    profile = str(plan.get("deck_profile") or "general")
    if profile == "consumer-retail" and slides:
        image_slides = [
            slide for slide in slides
            if slide.get("asset_ref") or slide.get("image_prompt")
        ]
        minimum = max(2, (len(slides) + 1) // 2)
        if len(image_slides) < minimum:
            issues.append(
                f"Image-led deck has {len(image_slides)} visual slides; it requires at least {minimum}."
            )
        if not (slides[0].get("asset_ref") or slides[0].get("image_prompt")):
            issues.append("Image-led deck cover requires a real or generated hero image.")
    for index, slide in enumerate(slides, start=1):
        if slide.get("type") == "image" and not (slide.get("asset_ref") or slide.get("image_prompt")):
            issues.append(f"Slide {index} requests an image layout without an image source.")
    return issues[:12]


def _render_presentation(
    plan: dict[str, Any],
    output: Path,
    theme_name: str = "mitsu_minimal",
    audience: str = "",
    base_presentation: Path | None = None,
    asset_lookup: dict[str, Path] | None = None,
    native_3d_assets: dict[str, Native3DAsset] | None = None,
    appearance: str = "dark",
    transition: str = "morph",
    include_speaker_notes: bool = False,
    theme_colors: dict[str, str] | None = None,
) -> None:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    colors = theme_colors or THEMES.get(theme_name, THEMES["mitsu_minimal"])
    rgb = lambda key: RGBColor.from_string(colors[key])
    prs = Presentation(base_presentation) if base_presentation else Presentation()
    if not base_presentation:
        prs.slide_width = Inches(13.333333)
        prs.slide_height = Inches(7.5)
    prs.core_properties.title = plan["title"]
    prs.core_properties.author = generated_by()
    prs.core_properties.subject = "Editable presentation created by MITSU"
    blank = next(
        (layout for layout in prs.slide_layouts if "blank" in str(layout.name).lower()),
        prs.slide_layouts[-1],
    )
    asset_lookup = asset_lookup or {}
    native_3d_assets = native_3d_assets or {}
    model_part_cache: dict[tuple[str, str], object] = {}
    deck_dark = str(appearance or "light").lower() == "dark"
    base_bg = "ink" if deck_dark else "paper"
    base_fg = "white" if deck_dark else "ink"
    secondary_fg = "muted"

    def rect(slide, x, y, w, h, fill, line=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        shape.line.color.rgb = rgb(line or fill)
        return shape

    def ring(slide, x, y, w, h, color="muted", width=0.6):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.background()
        shape.line.color.rgb = rgb(color)
        shape.line.width = Pt(width)
        return shape

    def text(slide, value, x, y, w, h, size, color, bold=False, font=BODY_FONT, align=PP_ALIGN.LEFT):
        if not str(value or "").strip():
            return None
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.text = str(value or "")
        paragraph.alignment = align
        paragraph.font.name = font
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = rgb(color)
        return box

    def bullet_list(slide, values, x, y, w, h, color="ink", size=18):
        values = values or []
        if not values:
            return
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = Inches(0.05)
        frame.margin_top = frame.margin_bottom = Inches(0.04)
        for index, value in enumerate(values):
            p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            p.text = f"—  {value}"
            p.font.name = BODY_FONT
            p.font.size = Pt(size)
            p.font.color.rgb = rgb(color)
            p.space_after = Pt(12)

    def labeled_list(slide, label, values, x, y, w, h, color=None, size=17):
        values = [str(item) for item in (values or []) if str(item).strip()]
        if not values:
            return None
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = Inches(0.03)
        frame.margin_top = frame.margin_bottom = Inches(0.04)
        first = frame.paragraphs[0]
        first.text = str(label or "").upper()
        first.font.name = DATA_FONT
        first.font.size = Pt(10)
        first.font.bold = True
        first.font.color.rgb = rgb("accent")
        first.space_after = Pt(12)
        for value in values:
            paragraph = frame.add_paragraph()
            paragraph.text = f"—  {value}"
            paragraph.font.name = BODY_FONT
            paragraph.font.size = Pt(size)
            paragraph.font.color.rgb = rgb(color or base_fg)
            paragraph.space_after = Pt(8)
        return box

    def canvas(slide, number, kicker, dark=False, source_note="", paint_background=True):
        bg = base_bg
        fg = base_fg
        if paint_background:
            rect(slide, 0, 0, 13.334, 7.5, bg)
        rect(slide, 0.55, 0.47, 0.075, 0.075, "accent")
        # Keep attribution visible without adding another floating text region.
        # A compact source line shares the small eyebrow area so normal slides
        # never exceed the four-region minimal-layout budget.
        eyebrow = str(kicker or "").strip()
        if source_note:
            eyebrow = f"{eyebrow}  •  {source_note}" if eyebrow else str(source_note)
        text(slide, eyebrow, 0.72, 0.34, 11.85, 0.32, 9, "accent", True, DATA_FONT)
        return fg

    def chart(slide, chart_spec, x=0.82, y=2.0, w=11.55, h=4.25):
        categories = chart_spec.get("categories") or []
        series = chart_spec.get("series") or []
        if not categories or not series:
            return False
        chart_data = CategoryChartData()
        chart_data.categories = categories
        for item in series:
            values = list(item.get("values") or [])
            values = (values + [0.0] * len(categories))[: len(categories)]
            chart_data.add_series(item.get("name") or "Series", values)
        chart_types = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "area": XL_CHART_TYPE.AREA,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
            "scatter": XL_CHART_TYPE.LINE_MARKERS,
        }
        graphic = slide.shapes.add_chart(
            chart_types.get(chart_spec.get("type"), XL_CHART_TYPE.COLUMN_CLUSTERED),
            Inches(x), Inches(y), Inches(w), Inches(h), chart_data,
        ).chart
        graphic.has_title = bool(chart_spec.get("title"))
        if graphic.has_title:
            graphic.chart_title.text_frame.text = chart_spec["title"]
        graphic.has_legend = len(series) > 1
        if graphic.has_legend:
            graphic.legend.position = XL_LEGEND_POSITION.BOTTOM
            graphic.legend.include_in_layout = False
        graphic.chart_style = 13
        try:
            graphic.chart_area.format.fill.solid()
            graphic.chart_area.format.fill.fore_color.rgb = rgb(base_bg)
            graphic.chart_area.format.line.fill.background()
            graphic.plot_area.format.fill.solid()
            graphic.plot_area.format.fill.fore_color.rgb = rgb(base_bg)
            palette = ("accent", "accent_2", "muted")
            for series_index, chart_series in enumerate(graphic.series):
                chart_series.format.fill.solid()
                chart_series.format.fill.fore_color.rgb = rgb(palette[series_index % len(palette)])
                chart_series.format.line.color.rgb = rgb(palette[series_index % len(palette)])
            if chart_spec.get("type") not in {"pie", "doughnut"}:
                graphic.value_axis.tick_labels.font.name = DATA_FONT
                graphic.value_axis.tick_labels.font.size = Pt(10)
                graphic.value_axis.tick_labels.font.color.rgb = rgb("muted")
                graphic.category_axis.tick_labels.font.name = BODY_FONT
                graphic.category_axis.tick_labels.font.size = Pt(10)
                graphic.category_axis.tick_labels.font.color.rgb = rgb(base_fg)
                graphic.value_axis.major_gridlines.format.line.color.rgb = rgb("muted")
        except Exception:
            # The editable chart is still valid if an Office version does not
            # expose one of the optional formatting properties.
            pass
        return True

    def table(slide, table_spec, x=0.8, y=2.0, w=11.7, h=4.35):
        headers = table_spec.get("headers") or []
        rows = table_spec.get("rows") or []
        if not headers or not rows:
            return False
        grid = slide.shapes.add_table(
            len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h),
        ).table
        for column in range(len(headers)):
            grid.columns[column].width = Inches(w / len(headers))
        values = [headers] + rows
        for row_index, row in enumerate(values):
            for column_index in range(len(headers)):
                cell = grid.cell(row_index, column_index)
                cell.text = str(row[column_index] if column_index < len(row) else "")
                cell.fill.solid()
                header_fill = "accent" if deck_dark else "ink"
                cell.fill.fore_color.rgb = rgb(header_fill if row_index == 0 else base_bg)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = BODY_FONT
                    paragraph.font.size = Pt(11 if row_index else 12)
                    paragraph.font.bold = row_index == 0
                    header_text = "ink" if deck_dark else "white"
                    paragraph.font.color.rgb = rgb(header_text if row_index == 0 else base_fg)
        return True

    def picture(slide, asset_ref, x=6.4, y=1.55, w=5.9, h=4.95):
        key = str(asset_ref or "").lower()
        native_asset = native_3d_assets.get(key)
        if native_asset:
            return add_native_3d_model(
                slide,
                native_asset,
                Inches(x),
                Inches(y),
                Inches(w),
                Inches(h),
                model_part_cache,
            )
        path = asset_lookup.get(key) or asset_lookup.get(Path(key).name)
        if not path or not Path(path).exists():
            return None
        shape = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
        try:
            from PIL import Image
            with Image.open(path) as source:
                source_ratio = source.width / max(1, source.height)
            box_ratio = w / h
            if source_ratio > box_ratio:
                crop = (1.0 - box_ratio / source_ratio) / 2.0
                shape.crop_left = crop
                shape.crop_right = crop
            elif source_ratio < box_ratio:
                crop = (1.0 - source_ratio / box_ratio) / 2.0
                shape.crop_top = crop
                shape.crop_bottom = crop
        except Exception:
            pass
        return shape

    def model_focus(slide, content, slide_number):
        """Reference-led composition: dominant 3D object plus three concise text regions."""
        key = str(content.get("asset_ref") or "").lower()
        if key not in native_3d_assets:
            return False
        native_asset = native_3d_assets[key]
        rect(slide, 0, 0, 13.334, 7.5, base_bg)
        ring(slide, -0.3, 0.38, 7.25, 6.72, "trace" if "trace" in colors else "muted", 0.4)
        if native_asset.poster_part:
            try:
                from io import BytesIO
                import zipfile

                with zipfile.ZipFile(native_asset.source_path) as source:
                    poster = slide.shapes.add_picture(
                        BytesIO(source.read(native_asset.poster_part)),
                        Inches(0.12), Inches(0.58), Inches(7.05), Inches(6.35),
                    )
                poster.name = f"!!MITSU 3D Poster — {native_asset.description or native_asset.name}"
            except Exception:
                pass
        model = picture(slide, content.get("asset_ref"), 0.12, 0.58, 7.05, 6.35)
        if model is None:
            return False

        label = str(content.get("kicker") or "").strip()
        eyebrow = f"{slide_number:02d}  ·  {label}" if label else f"{slide_number:02d}"
        text(slide, eyebrow, 7.78, 2.08, 4.75, 0.34, 9, "accent", True, DATA_FONT)
        title = text(
            slide, content.get("title"), 7.76, 2.52, 4.85, 1.62,
            39, base_fg, False, DISPLAY_FONT,
        )
        if title is not None:
            title.name = "model-focus-title"
        subtitle = text(
            slide, content.get("subtitle"), 7.8, 4.3, 4.65, 0.82,
            16, secondary_fg, False, BODY_FONT,
        )
        if subtitle is not None:
            subtitle.name = "model-focus-support"
        if content.get("source_note"):
            text(
                slide, content["source_note"], 7.8, 6.66, 4.65, 0.2,
                7, "muted", False, DATA_FONT,
            )
        return True

    for index, content in enumerate(plan["slides"], start=1):
        slide = prs.slides.add_slide(blank)
        kind = content["type"]
        asset_ref = content.get("asset_ref")
        model_focused = model_focus(slide, content, index)

        if model_focused:
            pass

        elif kind == "cover":
            rect(slide, 0, 0, 13.334, 7.5, base_bg)
            hero = picture(slide, asset_ref, 0.1, 0.45, 7.0, 6.6)
            text_x = 7.78 if hero is not None else 1.0
            text_w = 4.85 if hero is not None else 10.9
            text(slide, content["kicker"], text_x, 2.05, text_w, 0.34, 9, "accent", True, DATA_FONT)
            text(slide, content["title"], text_x, 2.5, text_w, 1.65, 40, base_fg, False, DISPLAY_FONT)
            subtitle = content["subtitle"] or plan.get("subtitle", "")
            text(slide, subtitle, text_x + 0.02, 4.3, text_w - 0.1, 0.82, 16, secondary_fg)

        elif kind == "image":
            variant = content.get("layout_variant") or "auto"
            if variant == "auto":
                variant = ("full_bleed", "image_right", "image_left")[index % 3]
            rect(slide, 0, 0, 13.334, 7.5, base_bg)
            if variant == "full_bleed":
                hero = picture(slide, asset_ref, 0, 0, 13.334, 7.5)
                if hero is not None:
                    rect(slide, 0, 0, 5.25, 7.5, base_bg)
                    canvas(slide, index, content["kicker"], deck_dark, content["source_note"], False)
                    text(slide, content["title"], 0.72, 1.25, 3.95, 1.65, 32, base_fg, True, DISPLAY_FONT)
                    text(slide, content["subtitle"], 0.75, 3.15, 3.85, 1.05, 18, secondary_fg)
                else:
                    canvas(slide, index, content["kicker"], deck_dark, content["source_note"], False)
                    text(slide, content["title"], 0.75, 1.25, 11.5, 1.55, 36, base_fg, True, DISPLAY_FONT)
                    text(slide, content["subtitle"], 0.78, 3.1, 9.8, 1.0, 20, secondary_fg)
            else:
                image_left = variant == "image_left"
                image_x = 0 if image_left else 7.15
                hero = picture(slide, asset_ref, image_x, 0, 6.18, 7.5)
                canvas(slide, index, content["kicker"], deck_dark, content["source_note"], False)
                text_x = 7.0 if image_left else 0.72
                text_w = 5.55
                text(slide, content["title"], text_x, 1.2, text_w, 1.55, 31, base_fg, True, DISPLAY_FONT)
                text(slide, content["subtitle"], text_x + 0.03, 3.02, text_w - 0.35, 1.0, 18, secondary_fg)
                if hero is None:
                    rect(slide, image_x + 0.5, 1.15, 0.11, 5.0, "accent")

        elif kind == "closing":
            rect(slide, 0, 0, 13.334, 7.5, base_bg)
            hero = picture(slide, asset_ref, 0.1, 0.45, 7.0, 6.6)
            text_x = 7.78 if hero is not None else 1.0
            text_w = 4.85 if hero is not None else 10.9
            text(slide, content["kicker"], text_x, 2.05, text_w, 0.34, 9, "accent", True, DATA_FONT)
            text(slide, content["title"], text_x, 2.5, text_w, 1.65, 40, base_fg, False, DISPLAY_FONT)
            text(slide, content["subtitle"], text_x + 0.02, 4.3, text_w - 0.1, 0.82, 16, secondary_fg)

        else:
            dark = deck_dark
            fg = canvas(slide, index, content["kicker"], dark, content["source_note"])

            if kind == "statement":
                text(slide, content["title"], 0.72, 1.05, 11.6, 1.65, 36, fg, True, DISPLAY_FONT)
                proof_items = content["visual_items"] or content["bullets"]
                proof_items = proof_items[:3]
                if proof_items:
                    labeled_list(
                        slide, content["visual_label"], proof_items,
                        0.75, 3.22, 11.1, 2.35, fg, 17,
                    )
                else:
                    text(slide, content["subtitle"], 0.75, 3.02, 10.4, 0.78, 19, secondary_fg)

            elif kind == "two_column":
                text(slide, content["title"], 0.72, 0.92, 11.3, 1.15, 31, fg, True, DISPLAY_FONT)
                left_items = content["bullets"]
                right_items = content["visual_items"]
                if right_items:
                    labeled_list(slide, content["left_label"], left_items, 0.75, 2.35, 5.25, 3.55, fg, 17)
                    rect(slide, 6.43, 2.25, 0.018, 3.62, "accent")
                    labeled_list(
                        slide, content["visual_label"], right_items[:4],
                        6.82, 2.35, 4.9, 3.2, fg, 16,
                    )
                else:
                    text(slide, content["subtitle"], 0.75, 2.45, 10.7, 0.8, 20, secondary_fg)
                    bullet_list(slide, left_items, 0.75, 3.55, 10.7, 2.2, fg, 18)

            elif kind in {"process", "timeline"}:
                text(slide, content["title"], 0.72, 0.92, 11.3, 1.15, 31, fg, True, DISPLAY_FONT)
                items = (content["visual_items"] or content["bullets"] or ["Frame", "Build", "Deliver"])[:4]
                box = labeled_list(
                    slide, "SEQUENCE", items,
                    0.75, 2.48, 11.1, 2.85, fg, 18,
                )
                if box is not None:
                    box.name = "process-label-01"
                elif content.get("subtitle"):
                    text(slide, content["subtitle"], 0.75, 3.0, 10.8, 0.65, 17, secondary_fg)

            elif kind == "comparison":
                text(slide, content["title"], 0.72, 0.92, 11.3, 1.15, 31, fg, True, DISPLAY_FONT)
                rect(slide, 6.43, 2.3, 0.018, 3.55, "accent")
                labeled_list(
                    slide, content["left_label"], content["left_items"] or content["bullets"],
                    0.78, 2.42, 5.2, 3.15, fg, 17,
                )
                labeled_list(
                    slide, content["right_label"], content["right_items"] or content["visual_items"],
                    6.82, 2.42, 5.2, 3.15, fg, 17,
                )

            elif kind == "chart":
                text(slide, content["title"], 0.72, 0.9, 11.3, 1.05, 30, fg, True, DISPLAY_FONT)
                if not chart(slide, content.get("chart", {}), 0.82, 2.05, 11.55, 4.45):
                    text(slide, content["subtitle"], 0.85, 2.45, 10.8, 0.8, 20, secondary_fg)
                    bullet_list(slide, content["bullets"] or content["visual_items"], 0.85, 3.55, 10.8, 2.2, fg, 18)

            elif kind in {"table", "appendix"}:
                text(slide, content["title"], 0.72, 0.9, 11.3, 1.05, 29, fg, True, DISPLAY_FONT)
                if not table(slide, content.get("table", {}), 0.8, 2.02, 11.7, 4.25):
                    text(slide, content["subtitle"], 0.85, 2.45, 10.8, 0.8, 19, secondary_fg)
                    bullet_list(slide, content["bullets"] or content["visual_items"], 0.85, 3.45, 10.8, 2.25, fg, 17)

        if include_speaker_notes and content.get("speaker_notes"):
            slide.notes_slide.notes_text_frame.text = content["speaker_notes"]
        if len(prs.slides) > 1:
            _apply_slide_transition(slide, transition)

    output.parent.mkdir(parents=True, exist_ok=True)
    temporary = output.with_name(f".{output.stem}.{uuid.uuid4().hex[:10]}.building.pptx")
    prs.save(temporary)
    temporary.replace(output)


def _open_presentation(path: Path) -> None:
    system = platform.system()
    if system == "Darwin":
        subprocess.Popen(["open", str(path)])
    elif system == "Windows":
        os.startfile(str(path))  # type: ignore[attr-defined]
    else:
        subprocess.Popen(["xdg-open", str(path)])


def _progress(callback: Callable | None, percent: int, phase: str, **extra: Any) -> None:
    if callback:
        callback(percent=percent, phase=phase, **extra)


def _check_cancel(cancel_flag: Any) -> None:
    if cancel_flag is not None and cancel_flag.is_set():
        raise RuntimeError("Presentation creation was cancelled.")


def _source_slide_count(path: Path | None) -> int:
    if not path:
        return 0
    from pptx import Presentation

    return len(Presentation(path).slides)


def _infer_theme_from_deck(path: Path | None) -> dict[str, str] | None:
    if not path:
        return None
    from collections import Counter
    from pptx import Presentation

    colors: Counter[str] = Counter()
    try:
        presentation = Presentation(path)
        for slide in list(presentation.slides)[:8]:
            for shape in slide.shapes:
                try:
                    rgb_value = shape.fill.fore_color.rgb
                    if rgb_value:
                        colors[str(rgb_value)] += 3
                except (AttributeError, TypeError):
                    pass
                if hasattr(shape, "text_frame") and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            try:
                                rgb_value = run.font.color.rgb
                                if rgb_value:
                                    colors[str(rgb_value)] += 1
                            except (AttributeError, TypeError):
                                pass
    except Exception:
        return None
    if not colors:
        return None

    ranked = [value for value, _ in colors.most_common(16)]

    def luminance(value: str) -> float:
        red, green, blue = int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)
        return 0.2126 * red + 0.7152 * green + 0.0722 * blue

    def saturation(value: str) -> int:
        channels = [int(value[index:index + 2], 16) for index in (0, 2, 4)]
        return max(channels) - min(channels)

    ink = min(ranked, key=luminance)
    paper = max(ranked, key=luminance)
    accent_candidates = [value for value in ranked if 45 < luminance(value) < 225]
    accent = max(accent_candidates or ranked, key=saturation)
    second_candidates = [value for value in accent_candidates if value != accent]
    accent_2 = max(second_candidates or ranked, key=lambda value: saturation(value) + abs(luminance(value) - luminance(accent)))
    muted_candidates = [value for value in ranked if value not in {ink, paper, accent, accent_2}]
    muted = min(muted_candidates or ranked, key=lambda value: abs(luminance(value) - 135))
    return {
        "ink": ink,
        "paper": paper,
        "accent": accent,
        "accent_2": accent_2,
        "muted": muted,
        "white": "FFFFFF",
    }


def _safe_output_for_source(output: Path, source: Path | None, mode: str) -> Path:
    if not source:
        return output
    try:
        same = output.resolve() == source.resolve()
    except OSError:
        same = str(output) == str(source)
    if not same:
        return output
    suffix = "edited" if mode == "edit" else "extended" if mode == "extend" else "redesigned"
    return output.with_name(f"{output.stem}_{suffix}.pptx")


def _apply_text_edits(source: Path, output: Path, edits: list[dict[str, Any]]) -> int:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    presentation = Presentation(source)
    applied = 0
    delete_numbers = sorted({
        int(edit.get("slide_number", 0))
        for edit in edits
        if edit.get("operation") == "delete_slide"
    }, reverse=True)
    for edit in edits:
        if edit.get("operation") == "delete_slide":
            continue
        slide_number = int(edit.get("slide_number", 0))
        if not 1 <= slide_number <= len(presentation.slides):
            continue
        find_text = str(edit.get("find_text") or "")
        replacement = str(edit.get("replacement_text") or "")
        if not find_text:
            continue
        slide = presentation.slides[slide_number - 1]
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                if find_text not in paragraph.text:
                    continue
                if len(paragraph.runs) > 1 and edit.get("operation") == "replace_text":
                    paragraph.text = paragraph.text.replace(find_text, replacement)
                runs = paragraph.runs
                if edit.get("operation") == "replace_text" and len(runs) == 1:
                    runs[0].text = runs[0].text.replace(find_text, replacement)
                for run in runs:
                    color = str(edit.get("font_color") or "")
                    if len(color) == 6:
                        run.font.color.rgb = RGBColor.from_string(color.upper())
                    try:
                        if edit.get("font_size") not in (None, ""):
                            run.font.size = Pt(float(edit["font_size"]))
                    except (TypeError, ValueError):
                        pass
                    if edit.get("bold") is not None:
                        run.font.bold = bool(edit["bold"])
                applied += 1
        notes = str(edit.get("notes") or "").strip()
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    for slide_number in delete_numbers:
        if 1 <= slide_number <= len(presentation.slides):
            slide_id = presentation.slides._sldIdLst[slide_number - 1]
            presentation.part.drop_rel(slide_id.rId)
            presentation.slides._sldIdLst.remove(slide_id)
            applied += 1
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)
    return applied


def _revise_plan(
    plan: dict[str, Any],
    critique: dict[str, Any],
    request: PresentationRequest,
) -> dict[str, Any]:
    from google import genai

    prompt = (
        "You are revising a PowerPoint plan after rendered QA. Fix only the reported "
        "problems while preserving exact facts, citations, slide count, and JSON shape. "
        "Shorten text when clipping or crowding is reported. Return JSON only.\n\n"
        f"REQUEST: {request.topic}\n\nCRITIQUE:\n{json.dumps(critique, ensure_ascii=False)}\n\n"
        f"CURRENT PLAN:\n{json.dumps(plan, ensure_ascii=False)}"
    )
    client = genai.Client(api_key=_api_key())
    response = client.models.generate_content(
        model=DEFAULT_MODEL,
        contents=prompt,
        config={"response_mime_type": "application/json"},
    )
    revised = _clean_json(response.text)
    return _normalize_plan(
        revised,
        request.topic,
        len(plan.get("slides", [])),
        bookends=request.resolved_mode() in {"create", "redesign"},
    )


def build_presentation(
    parameters: dict,
    *,
    player=None,
    cancel_flag: Any = None,
    progress_callback: Callable | None = None,
) -> PresentationResult:
    """Build and verify a presentation synchronously.

    Live requests enqueue this function as a specialized background job. Direct
    callers retain a deterministic synchronous entry point for testing and agent
    tasks that are already running in the shared queue.
    """
    request = PresentationRequest.from_parameters(parameters)
    mode = request.resolved_mode()
    source_deck = request.source_deck()
    if mode in {"edit", "extend"} and not source_deck:
        raise ValueError(f"Presentation mode '{mode}' requires an existing PPTX source or template.")
    if source_deck and (not source_deck.exists() or source_deck.suffix.lower() != ".pptx"):
        raise FileNotFoundError(f"PowerPoint source deck was not found: {source_deck}")

    source_count = _source_slide_count(source_deck)
    theme_colors = _infer_theme_from_deck(source_deck) if source_deck and not request.theme_explicit else None
    if mode == "redesign" and source_count and not request.slide_count_explicit:
        request.slide_count = min(MAX_SLIDES, max(MIN_SLIDES, source_count))
    plan_count = request.slide_count
    if mode == "extend":
        plan_count = (
            max(1, request.slide_count - source_count)
            if request.slide_count_explicit and request.slide_count > source_count
            else min(3, MAX_SLIDES - source_count)
        )
        if plan_count < 1:
            raise ValueError("The source deck already meets or exceeds the requested 50-slide limit.")

    output = _safe_output_for_source(_output_path(request.output_path, request.title), source_deck, mode)
    output.parent.mkdir(parents=True, exist_ok=True)
    workspace = BASE_DIR / "outputs" / "presentation_jobs" / uuid.uuid4().hex[:12]
    workspace.mkdir(parents=True, exist_ok=True)
    output_lock = _output_lock_for(output)
    output_lock.acquire()
    warnings: list[str] = []
    pdf_path: Path | None = None
    qa: dict[str, Any] = {}

    try:
        _progress(progress_callback, 3, "Validating presentation request")
        _check_cancel(cancel_flag)
        bundle = ingest_sources(request, progress_callback, workspace)
        warnings.extend(bundle.warnings)
        (workspace / "source-manifest.json").write_text(source_manifest(bundle), encoding="utf-8")
        pptx_sources = [
            Path(record.location)
            for record in bundle.records
            if record.kind == "pptx" and Path(record.location).is_file()
        ]
        native_assets = native_3d_lookup(pptx_sources)
        unique_native_assets = {
            asset.asset_ref: asset for asset in native_assets.values()
        }
        qa["native_3d_inventory"] = {
            "requested": request.use_native_3d,
            "available": len(unique_native_assets),
            "assets": [asset.metadata() for asset in unique_native_assets.values()],
        }
        if request.use_native_3d and not unique_native_assets:
            raise ValueError(
                "Native PowerPoint 3D was requested, but the supplied files contain no embedded "
                "3D models. Attach a PPTX with the model, or continue without 3D."
            )

        _progress(progress_callback, 24, "Planning narrative and evidence")
        _check_cancel(cancel_flag)
        ai_quota_exhausted = False
        if _presentation_quota_cooldown_active() and mode in {"create", "redesign"}:
            ai_quota_exhausted = True
            warnings.append(
                "Gemini presentation planning is in a short quota cooldown. MITSU completed "
                "a minimal local deck without making another API request."
            )
            qa["quota_fallback"] = {
                "used": True,
                "stage": "planning cooldown",
                "additional_ai_calls": 0,
            }
            raw_plan = _quota_fallback_plan(request, plan_count, bundle.text)
        else:
            try:
                raw_plan = _plan_presentation(
                    request.topic,
                    request.title,
                    request.audience,
                    plan_count,
                    request.tone,
                    bundle.text[:220_000],
                    request=request,
                    bundle=bundle,
                    cancel_flag=cancel_flag,
                )
            except Exception as exc:
                if not _is_quota_error(exc) or mode not in {"create", "redesign"}:
                    raise
                _mark_presentation_quota_cooldown()
                ai_quota_exhausted = True
                warnings.append(
                    "Gemini planning quota was unavailable. MITSU completed a minimal local deck "
                    "instead of failing or making more API requests."
                )
                qa["quota_fallback"] = {
                    "used": True,
                    "stage": "planning",
                    "additional_ai_calls": 0,
                }
                raw_plan = _quota_fallback_plan(request, plan_count, bundle.text)
        plan = (
            _normalize_edit_plan(raw_plan, request.topic)
            if mode == "edit"
            else _normalize_plan(
                raw_plan,
                request.topic,
                plan_count,
                bookends=mode in {"create", "redesign"},
            )
        )
        if mode != "edit" and request.quality != "fast" and not ai_quota_exhausted:
            plan = _enforce_visual_story(plan, request.topic)
        if mode != "edit":
            plan = _enforce_native_3d_usage(plan, native_assets, request.use_native_3d)
        warnings.extend(bundle.warnings)

        if mode != "edit":
            preflight_issues = _plan_quality_issues(plan)
            qa["plan_preflight"] = {
                "passed": not preflight_issues,
                "issues": preflight_issues,
            }
            if preflight_issues:
                warnings.append(
                    "Local presentation preflight retained a few non-blocking story notes: "
                    + "; ".join(preflight_issues[:3])
                )

        _progress(progress_callback, 43, "Preparing visual assets")
        _check_cancel(cancel_flag)
        needs_generated_images = any(
            str(slide.get("image_prompt") or "").strip() and not slide.get("asset_ref")
            for slide in plan.get("slides", [])
        )
        asset_lookup, asset_warnings = generate_planned_assets(
            plan,
            request,
            bundle,
            workspace,
            progress_callback,
            cancel_flag,
            _api_key() if needs_generated_images and not ai_quota_exhausted else "",
        )
        warnings.extend(asset_warnings)
        if any("quota" in warning.lower() for warning in asset_warnings):
            ai_quota_exhausted = True

        _progress(progress_callback, 58, "Building editable PowerPoint")
        _check_cancel(cancel_flag)
        if mode == "edit":
            edited_base = workspace / "edited-base.pptx"
            applied = _apply_text_edits(source_deck, edited_base, plan.get("edits", []))
            if not applied and not plan.get("slides"):
                warnings.append("No matching editable text was found in the source deck.")
            if plan.get("slides"):
                _render_presentation(
                    plan, output, request.theme, request.audience,
                    base_presentation=edited_base,
                    asset_lookup=asset_lookup,
                    native_3d_assets=native_assets,
                    appearance=request.resolved_appearance(),
                    transition=request.transition,
                    include_speaker_notes=request.include_speaker_notes,
                    theme_colors=theme_colors,
                )
            else:
                shutil.copy2(edited_base, output)
        else:
            base = source_deck if mode == "extend" else None
            _render_presentation(
                plan, output, request.theme, request.audience,
                base_presentation=base,
                asset_lookup=asset_lookup,
                native_3d_assets=native_assets,
                appearance=request.resolved_appearance(),
                transition=request.transition,
                include_speaker_notes=request.include_speaker_notes,
                theme_colors=theme_colors,
            )

        _progress(progress_callback, 70, "Verifying PowerPoint package")
        _check_cancel(cancel_flag)
        deleted_slides = len({
            int(edit.get("slide_number", 0))
            for edit in plan.get("edits", [])
            if edit.get("operation") == "delete_slide"
            and 1 <= int(edit.get("slide_number", 0)) <= source_count
        }) if mode == "edit" else 0
        expected = (
            source_count - deleted_slides + len(plan.get("slides", []))
            if mode in {"edit", "extend"}
            else len(plan["slides"])
        )
        qa["package"] = verify_pptx(output, expected)
        if request.use_native_3d and qa["package"]["native_3d_shape_count"] < 1:
            raise RuntimeError(
                "The presentation requested native 3D, but verification found no 3D model shapes."
            )
        expected_new_transitions = (
            0
            if not plan.get("slides")
            else max(0, len(plan.get("slides", [])) - (0 if source_count else 1))
        )
        if request.transition == "morph" and (
            qa["package"]["morph_slide_count"] < expected_new_transitions
        ):
            warnings.append(
                "Some PowerPoint versions retained the Fade compatibility transition instead "
                "of exposing every Morph marker. The presentation was kept rather than discarded."
            )
            qa["transition_fallback"] = "fade"
        elif request.transition == "fade" and (
            qa["package"]["transition_slide_count"] < expected_new_transitions
        ):
            warnings.append(
                "PowerPoint did not expose every requested Fade transition, but the verified "
                "presentation remains usable."
            )

        if request.export_pdf:
            _progress(progress_callback, 76, "Exporting presentation PDF")
            pdf_candidate = output.with_suffix(".pdf")
            pdf_path, export_warning = export_pdf(output, pdf_candidate)
            if export_warning and export_warning not in warnings:
                warnings.append(export_warning)

        # Standard quality uses local render/package QA only. The optional
        # premium path gets one art-director pass and at most one repair call.
        # This replaces the previous multi-round critique loop that could burn
        # several Gemini requests on a single deck.
        if pdf_path and request.quality != "fast":
            _progress(progress_callback, 82, "Rendering slides for visual QA")
            preview_dir = workspace / "preview-final"
            renders, render_warning = render_pdf(pdf_path, preview_dir)
            if render_warning:
                warnings.append(render_warning)
            else:
                qa["local_render"] = {
                    "rendered_slide_count": len(renders),
                    "expected_slide_count": expected,
                    "passed": len(renders) == expected,
                }
                contact_sheet = create_contact_sheet(renders, preview_dir / "contact-sheet.png")
                if request.quality == "premium" and mode != "edit" and not ai_quota_exhausted:
                    critique, critique_warning = critique_renders(
                        renders, contact_sheet, request.topic, _api_key(),
                    )
                    if critique_warning:
                        warnings.append(critique_warning)
                        if "quota" in critique_warning.lower():
                            _mark_presentation_quota_cooldown()
                            ai_quota_exhausted = True
                    else:
                        qa["premium_visual_review"] = critique
                        blocking = critique.get("blocking_issues") if isinstance(critique, dict) else []
                        score = int(critique.get("score", 0) or 0) if isinstance(critique, dict) else 0
                        if blocking or score < 80:
                            _progress(progress_callback, 88, "Applying one premium visual repair")
                            _check_cancel(cancel_flag)
                            try:
                                plan = _revise_plan(plan, critique, request)
                            except Exception as exc:
                                if _is_quota_error(exc):
                                    _mark_presentation_quota_cooldown()
                                    ai_quota_exhausted = True
                                    warnings.append(
                                        "Gemini quota was reached during optional premium repair. "
                                        "MITSU kept the verified PowerPoint instead of failing."
                                    )
                                else:
                                    warnings.append(f"Premium visual repair was unavailable: {exc}")
                            else:
                                plan = _enforce_visual_story(plan, request.topic)
                                plan = _enforce_native_3d_usage(plan, native_assets, request.use_native_3d)
                                _render_presentation(
                                    plan, output, request.theme, request.audience,
                                    base_presentation=source_deck if mode == "extend" else None,
                                    asset_lookup=asset_lookup,
                                    native_3d_assets=native_assets,
                                    appearance=request.resolved_appearance(),
                                    transition=request.transition,
                                    include_speaker_notes=request.include_speaker_notes,
                                    theme_colors=theme_colors,
                                )
                                qa["package"] = verify_pptx(output, expected)
                                repaired_pdf, export_warning = export_pdf(output, output.with_suffix(".pdf"))
                                if repaired_pdf:
                                    pdf_path = repaired_pdf
                                if export_warning and export_warning not in warnings:
                                    warnings.append(export_warning)

        mark_created_file(output, f"PowerPoint presentation about {request.topic}")
        if pdf_path:
            mark_created_file(pdf_path, f"PDF export of PowerPoint presentation about {request.topic}")
        if request.open_after_create:
            _open_presentation(output)
        if player and hasattr(player, "write_log"):
            player.write_log(f"PRESENTATION: Created {output.name}")
        _progress(
            progress_callback, 100, "Presentation complete",
            artifacts=[str(output)] + ([str(pdf_path)] if pdf_path else []),
            warnings=warnings,
        )
        return PresentationResult(
            pptx_path=output,
            pdf_path=pdf_path,
            slide_count=qa["package"]["slide_count"],
            mode=mode,
            quality=request.quality,
            warnings=warnings,
            qa=qa,
        )
    finally:
        clean_preview_dir(workspace)
        output_lock.release()


def create_presentation(parameters: dict, response=None, player=None, session_memory=None) -> str:
    """Generate an editable PowerPoint deck from a spoken or typed request.

    ``response`` and ``session_memory`` are accepted for compatibility with the
    rest of MITSU' action interface.  Presentation creation is intentionally a
    dedicated action so requests never fall through to generic code generation
    or desktop automation.
    """
    try:
        if player and hasattr(player, "write_log"):
            player.write_log("PRESENTATION: Starting universal presentation pipeline...")
        return str(build_presentation(parameters or {}, player=player))
    except Exception as exc:
        return f"Could not create the PowerPoint presentation: {exc}"


def queue_presentation(
    parameters: dict,
    player=None,
    speak=None,
    *,
    visible: bool = False,
) -> str:
    """Submit a presentation as a cancellable, progress-reporting job."""
    request = PresentationRequest.from_parameters(parameters)
    from agent.task_queue import TaskPriority, get_queue

    def runner(cancel_flag, progress_callback):
        last_phase = {"value": ""}

        def report(**update):
            progress_callback(**update)
            phase = str(update.get("phase") or "")
            if player and hasattr(player, "update_presentation_progress"):
                player.update_presentation_progress(
                    title=request.title,
                    percent=update.get("percent", 0),
                    phase=phase or "Building presentation",
                    visible=visible,
                    artifacts=update.get("artifacts") or [],
                    warnings=update.get("warnings") or [],
                )
            if phase and phase != last_phase["value"] and player and hasattr(player, "write_log"):
                player.write_log(f"PRESENTATION [{update.get('percent', 0)}%]: {phase}")
                last_phase["value"] = phase

        try:
            return build_presentation(
                parameters,
                player=player,
                cancel_flag=cancel_flag,
                progress_callback=report,
            )
        except Exception as exc:
            if player and hasattr(player, "finish_presentation_progress"):
                state = "cancelled" if cancel_flag.is_set() else "failed"
                player.finish_presentation_progress(state, str(exc)[:160])
            raise

    if player and hasattr(player, "show_presentation_progress"):
        player.show_presentation_progress(request.title, visible=visible)
    try:
        on_cancel = None
        if player and hasattr(player, "finish_presentation_progress"):
            on_cancel = lambda: player.finish_presentation_progress(
                "cancelled", "Presentation task cancelled"
            )
        task_id = get_queue().submit_job(
            goal=f"Create presentation: {request.title}",
            runner=runner,
            priority=TaskPriority.NORMAL,
            speak=speak,
            on_cancel=on_cancel,
            kind="presentation",
        )
    except Exception as exc:
        if player and hasattr(player, "finish_presentation_progress"):
            player.finish_presentation_progress("failed", str(exc)[:160])
        raise
    if visible:
        return (
            f"Presentation task is visible (ID: {task_id}). "
            "I will show each real build phase while I work and announce when the files are ready."
        )
    return (
        f"Presentation task started in the background (ID: {task_id}). "
        "A compact status indicator will remain available, and I will announce when the files are ready."
    )


def _presentation_execution_mode(parameters: dict | None) -> str:
    value = str((parameters or {}).get("execution_mode", "ask") or "ask").strip().lower()
    aliases = {
        "foreground": "visible",
        "live": "visible",
        "show": "visible",
        "watch": "visible",
        "quiet": "background",
    }
    value = aliases.get(value, value)
    if value not in {"ask", "background", "visible"}:
        raise ValueError("Presentation execution mode must be ask, background, or visible.")
    return value


def _presentation_3d_preference(parameters: dict | None) -> bool | None:
    params = parameters or {}
    choice = str(params.get("three_d_mode") or "").strip().lower()
    if choice in {"yes", "use", "include", "required"}:
        return True
    if choice in {"no", "skip", "none", "without"}:
        return False
    if "use_native_3d" in params:
        raw = params.get("use_native_3d")
        if isinstance(raw, bool):
            return raw
        value = str(raw or "").strip().lower()
        if value in {"true", "yes", "on", "1"}:
            return True
        if value in {"false", "no", "off", "0"}:
            return False
    topic = str(
        params.get("topic") or params.get("instructions") or params.get("description") or ""
    ).lower()
    if re.search(r"\b(?:without|no|skip|do not use)\b.{0,24}\b3d\b", topic):
        return False
    if re.search(r"\b(?:use|include|with|add)\b.{0,24}\b3d\b", topic):
        return True
    return None


def _remember_pending_presentation(parameters: dict) -> None:
    global _pending_presentation_parameters, _pending_presentation_created_at
    stored = dict(parameters)
    stored.pop("execution_mode", None)
    from core.tenant import get_current_user_id

    user_id = get_current_user_id()
    with _pending_presentation_lock:
        if user_id:
            _tenant_pending_presentations[user_id] = (stored, time.monotonic())
        else:
            _pending_presentation_parameters = stored
            _pending_presentation_created_at = time.monotonic()


def _take_pending_presentation(parameters: dict) -> dict | None:
    global _pending_presentation_parameters, _pending_presentation_created_at
    supplied = dict(parameters)
    supplied.pop("execution_mode", None)
    from core.tenant import get_current_user_id

    user_id = get_current_user_id()
    with _pending_presentation_lock:
        if user_id:
            pending, created_at = _tenant_pending_presentations.pop(user_id, (None, 0.0))
            age = time.monotonic() - created_at
        else:
            pending = _pending_presentation_parameters
            age = time.monotonic() - _pending_presentation_created_at
            _pending_presentation_parameters = None
            _pending_presentation_created_at = 0.0
    if pending is not None and age <= PRESENTATION_PENDING_TTL_SECONDS:
        merged = dict(pending)
        merged.update({
            key: value for key, value in supplied.items()
            if value not in (None, "", [])
        })
        return merged
    return None


def request_presentation(parameters: dict, player=None, speak=None) -> str:
    """Collect 3D and run-mode choices, then start the confirmed task."""
    params = dict(parameters or {})
    mode = _presentation_execution_mode(params)
    has_brief = bool(
        params.get("topic") or params.get("instructions") or params.get("description")
    )
    if has_brief:
        PresentationRequest.from_parameters(params)
        resolved = params
    else:
        resolved = _take_pending_presentation(params)

    if resolved is None:
        return (
            "Tell me which presentation you want first. I will ask whether you want a 3D model, "
            "then whether you want to see the task or keep it in the background."
        )

    three_d_preference = _presentation_3d_preference(resolved)
    if three_d_preference is None:
        _remember_pending_presentation(resolved)
        if player and hasattr(player, "write_log"):
            player.write_log(
                "PRESENTATION: Waiting for native 3D preference"
            )
        return PRESENTATION_3D_QUESTION

    resolved["use_native_3d"] = three_d_preference
    if resolved.pop("_awaiting_3d_source", False):
        supplied = resolved.get("source_files")
        supplied_paths = list(supplied) if isinstance(supplied, list) else ([supplied] if supplied else [])
        for raw in supplied_paths:
            path = Path(str(raw)).expanduser()
            if path.suffix.lower() == ".pptx" and inspect_native_3d(path):
                resolved["model_source_file"] = str(path)
                remaining = [item for item in supplied_paths if str(item) != str(raw)]
                if remaining:
                    resolved["source_files"] = remaining
                else:
                    resolved.pop("source_files", None)
                break
    if three_d_preference and not _request_native_3d_assets(resolved):
        resolved["_awaiting_3d_source"] = True
        _remember_pending_presentation(resolved)
        if player and hasattr(player, "write_log"):
            player.write_log("PRESENTATION: Waiting for a PPTX containing a native 3D model")
        return PRESENTATION_3D_SOURCE_QUESTION
    if mode == "ask":
        request = PresentationRequest.from_parameters(resolved)
        _remember_pending_presentation(resolved)
        if player and hasattr(player, "write_log"):
            player.write_log(
                f"PRESENTATION: Waiting for run-mode choice for {request.title[:70]}"
            )
        return PRESENTATION_RUN_MODE_QUESTION

    return queue_presentation(
        resolved,
        player=player,
        speak=speak,
        visible=(mode == "visible"),
    )


def presentation_maker(parameters: dict, response=None, player=None, session_memory=None) -> str:
    """Backward-compatible alias for the dedicated presentation action."""
    return create_presentation(
        parameters=parameters,
        response=response,
        player=player,
        session_memory=session_memory,
    )
