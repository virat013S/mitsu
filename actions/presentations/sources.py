"""Source ingestion and provenance tracking for presentation jobs."""

from __future__ import annotations

import csv
import hashlib
import json
import mimetypes
from pathlib import Path
from typing import Callable

from .models import PresentationRequest, SourceBundle, SourceRecord
from .models3d import inspect_native_3d


TEXT_EXTENSIONS = {
    ".txt", ".md", ".rst", ".json", ".jsonl", ".xml", ".yaml", ".yml",
    ".py", ".js", ".ts", ".html", ".css", ".sql",
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff"}
AUDIO_EXTENSIONS = {".wav", ".mp3", ".aiff", ".aac", ".ogg", ".flac", ".m4a"}
VIDEO_EXTENSIONS = {".mp4", ".mpeg", ".mov", ".avi", ".webm", ".wmv", ".mpg", ".3gp"}
SUPPORTED_EXTENSIONS = (
    TEXT_EXTENSIONS
    | IMAGE_EXTENSIONS
    | AUDIO_EXTENSIONS
    | VIDEO_EXTENSIONS
    | {".csv", ".xlsx", ".xls", ".docx", ".pptx", ".pdf"}
)


def _progress(callback: Callable | None, percent: int, phase: str) -> None:
    if callback:
        callback(percent=percent, phase=phase)


def _read_text(path: Path, maximum: int = 80_000) -> str:
    return path.read_text(encoding="utf-8", errors="replace")[:maximum]


def _read_csv(path: Path, maximum_rows: int = 250) -> tuple[str, dict]:
    with path.open("r", encoding="utf-8", errors="replace", newline="") as handle:
        reader = csv.reader(handle)
        rows = []
        for index, row in enumerate(reader):
            if index >= maximum_rows:
                break
            rows.append(row)
    preview = "\n".join(", ".join(str(cell) for cell in row) for row in rows)
    return preview[:80_000], {"rows_loaded": len(rows)}


def _read_xlsx(path: Path, maximum_rows: int = 250) -> tuple[str, dict]:
    if path.suffix.lower() == ".xls":
        try:
            import pandas as pd
        except ImportError as exc:
            raise RuntimeError("Legacy Excel sources require pandas and xlrd.") from exc
        workbook = pd.read_excel(path, sheet_name=None, nrows=maximum_rows)
        blocks = []
        for name, frame in workbook.items():
            blocks.append(f"SHEET: {name}\n{frame.to_string(index=False)}")
        return "\n\n".join(blocks)[:120_000], {"sheets": list(workbook)}
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("Excel sources require openpyxl.") from exc
    workbook = load_workbook(path, read_only=True, data_only=True)
    blocks = []
    metadata = {"sheets": workbook.sheetnames}
    for sheet in workbook.worksheets:
        lines = [f"SHEET: {sheet.title}"]
        for index, row in enumerate(sheet.iter_rows(values_only=True)):
            if index >= maximum_rows:
                break
            lines.append(" | ".join("" if value is None else str(value) for value in row))
        blocks.append("\n".join(lines))
    workbook.close()
    return "\n\n".join(blocks)[:120_000], metadata


def _read_docx(path: Path) -> tuple[str, dict]:
    from docx import Document

    document = Document(path)
    chunks = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            chunks.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(chunks)[:120_000], {
        "paragraphs": len(document.paragraphs),
        "tables": len(document.tables),
    }


def _read_pptx(path: Path) -> tuple[str, dict]:
    from pptx import Presentation

    presentation = Presentation(path)
    chunks = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_text = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        chunks.append(f"SLIDE {index}\n" + "\n".join(slide_text))
    native_models = inspect_native_3d(path)
    if native_models:
        chunks.append(
            "NATIVE POWERPOINT 3D ASSETS (editable and rotatable in PowerPoint)\n"
            + "\n".join(
                f"- asset_ref={asset.asset_ref} | name={asset.name} | "
                f"description={asset.description or 'unspecified'} | source_slide={asset.source_slide}"
                for asset in native_models
            )
        )
    return "\n\n".join(chunks)[:150_000], {
        "slides": len(presentation.slides),
        "width": presentation.slide_width,
        "height": presentation.slide_height,
        "native_3d_count": len(native_models),
        "native_3d_assets": [asset.metadata() for asset in native_models],
    }


def _read_pdf(path: Path) -> tuple[str, dict]:
    try:
        import fitz

        document = fitz.open(path)
        text = "\n".join(page.get_text("text") for page in document)
        pages = len(document)
        document.close()
        return text[:150_000], {"pages": pages, "native_visual_source": True}
    except ImportError:
        return "", {"native_visual_source": True}


def _record_for_path(path: Path) -> SourceRecord:
    if not path.exists() or not path.is_file():
        raise FileNotFoundError(f"Presentation source file was not found: {path}")
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported presentation source type: {suffix or path.name}")

    text = ""
    metadata: dict = {"bytes": path.stat().st_size, "mime_type": mimetypes.guess_type(path)[0] or ""}
    media_path: Path | None = None
    kind = suffix.lstrip(".") or "file"

    if suffix in TEXT_EXTENSIONS:
        text = _read_text(path)
    elif suffix == ".csv":
        text, extra = _read_csv(path)
        metadata.update(extra)
    elif suffix in {".xlsx", ".xls"}:
        text, extra = _read_xlsx(path)
        metadata.update(extra)
    elif suffix == ".docx":
        text, extra = _read_docx(path)
        metadata.update(extra)
    elif suffix == ".pptx":
        text, extra = _read_pptx(path)
        metadata.update(extra)
    elif suffix == ".pdf":
        text, extra = _read_pdf(path)
        metadata.update(extra)
        media_path = path
    else:
        media_path = path

    return SourceRecord(
        label=path.name,
        kind=kind,
        location=str(path.resolve()),
        text=text,
        media_path=media_path,
        provenance="user-provided",
        metadata=metadata,
    )


def _read_url(url: str, workspace: Path | None = None) -> SourceRecord:
    if not url.lower().startswith(("http://", "https://")):
        raise ValueError(f"Unsupported source URL: {url}")
    try:
        import requests
        from bs4 import BeautifulSoup

        response = requests.get(url, timeout=20, headers={"User-Agent": "MITSU-Presentation/1.0"})
        response.raise_for_status()
        content_type = response.headers.get("content-type", "")
        if "text/html" in content_type:
            soup = BeautifulSoup(response.text, "html.parser")
            for element in soup(["script", "style", "nav", "footer"]):
                element.decompose()
            text = "\n".join(line.strip() for line in soup.get_text("\n").splitlines() if line.strip())
            media_path = None
            kind = "url"
        elif any(token in content_type for token in ("pdf", "image/", "audio/", "video/")):
            if workspace is None:
                raise RuntimeError("A workspace is required for binary URL sources.")
            suffix = mimetypes.guess_extension(content_type.split(";", 1)[0].strip()) or Path(url).suffix or ".bin"
            source_dir = workspace / "sources"
            source_dir.mkdir(parents=True, exist_ok=True)
            media_path = source_dir / f"url-{hashlib.sha256(url.encode()).hexdigest()[:12]}{suffix}"
            media_path.write_bytes(response.content)
            text = ""
            kind = suffix.lstrip(".") or "url-media"
        else:
            text = response.text
            media_path = None
            kind = "url"
        return SourceRecord(
            label=url,
            kind=kind,
            location=url,
            text=text[:100_000],
            media_path=media_path,
            provenance="user-provided URL",
            metadata={"content_type": content_type},
        )
    except Exception as exc:
        raise RuntimeError(f"Could not retrieve presentation source URL {url}: {exc}") from exc


def ingest_sources(
    request: PresentationRequest,
    progress_callback: Callable | None = None,
    workspace: Path | None = None,
) -> SourceBundle:
    bundle = SourceBundle()
    paths = list(request.source_files)
    if request.template_file and request.template_file not in paths:
        paths.append(request.template_file)
    if request.model_source_file and request.model_source_file not in paths:
        paths.append(request.model_source_file)
    total = max(1, len(paths) + len(request.source_urls))
    completed = 0

    for path in paths:
        record = _record_for_path(path)
        if request.model_source_file and path == request.model_source_file:
            # This deck is a model library, not narrative input. Keep its file
            # location and 3D inventory available to the package builder while
            # preventing its slide copy from leaking into the new storyline.
            record.text = ""
            record.provenance = "user-provided PowerPoint 3D model library"
            record.metadata["model_library_only"] = True
        bundle.records.append(record)
        completed += 1
        _progress(progress_callback, 5 + int(15 * completed / total), f"Reading {path.name}")

    for url in request.source_urls:
        bundle.records.append(_read_url(url, workspace))
        completed += 1
        _progress(progress_callback, 5 + int(15 * completed / total), "Reading source URL")

    if request.allow_web_research:
        try:
            from actions.web_search import web_search

            research = web_search({"query": request.topic, "mode": "search"}, player=None)
            if research:
                bundle.records.append(SourceRecord(
                    label="MITSU web research",
                    kind="research",
                    location=f"web search: {request.topic}",
                    text=str(research)[:80_000],
                    provenance="web research authorized by user",
                ))
        except Exception as exc:
            bundle.warnings.append(f"Web research was unavailable: {exc}")

    return bundle


def source_manifest(bundle: SourceBundle) -> str:
    payload = []
    for index, record in enumerate(bundle.records, start=1):
        payload.append({
            "id": index,
            "label": record.label,
            "kind": record.kind,
            "location": record.location,
            "provenance": record.provenance,
            "metadata": record.metadata,
        })
    return json.dumps(payload, indent=2, ensure_ascii=False)
